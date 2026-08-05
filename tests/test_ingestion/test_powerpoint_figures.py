"""PowerPoint slide text, chart, notes, and positioned-picture ingestion."""
from __future__ import annotations

import io
from pathlib import Path
from unittest.mock import patch

from PIL import Image


def _png_bytes(color=(38, 82, 150)) -> bytes:
    image = Image.new("RGB", (640, 360), color)
    stream = io.BytesIO()
    image.save(stream, format="PNG")
    return stream.getvalue()


def _make_pptx(path: Path, *, repeated_slide: bool = False) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    image_path = path.parent / "_pptx_figure.png"
    image_path.write_bytes(_png_bytes())
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Control Plane"
    before = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(5), Inches(0.5))
    before.text = "Traffic enters through EDGE-01."
    picture = slide.shapes.add_picture(
        str(image_path), Inches(1), Inches(2), width=Inches(4), height=Inches(2.25),
    )
    picture._element.xpath(".//p:cNvPr")[0].set(
        "descr", "Redundant controller topology",
    )
    caption = slide.shapes.add_textbox(Inches(1), Inches(4.35), Inches(4), Inches(0.4))
    caption.text = "Figure 2. Controller topology"
    after = slide.shapes.add_textbox(Inches(0.6), Inches(5), Inches(6), Inches(0.5))
    after.text = "Controllers run in separate availability zones."
    slide.notes_slide.notes_text_frame.text = "Failover uses a three-node quorum."

    if repeated_slide:
        second = presentation.slides.add_slide(presentation.slide_layouts[5])
        second.shapes.title.text = "Recovery"
        second.shapes.add_picture(
            str(image_path), Inches(1), Inches(2), width=Inches(4), height=Inches(2.25),
        )

    presentation.save(path)
    image_path.unlink(missing_ok=True)


def test_pptx_parser_preserves_slide_picture_position_and_context(tmp_path):
    from src.ingestion.parser import parse_document

    path = tmp_path / "architecture.pptx"
    _make_pptx(path)
    parsed = parse_document(path)

    assert parsed.doc_type == "pptx"
    assert "## Slide 1: Control Plane" in parsed.text
    assert "Speaker notes:\nFailover uses a three-node quorum." in parsed.text
    figure_index = next(
        i for i, block in enumerate(parsed.blocks) if block.block_type == "figure"
    )
    placement = parsed.blocks[figure_index].figure
    assert placement is not None
    assert placement.figure_id == "s1-fig-001"
    assert placement.page == 0
    assert placement.bbox is not None
    assert placement.section_path == ["Control Plane"]
    assert placement.caption == "Figure 2. Controller topology"
    assert placement.alt_text == "Redundant controller topology"
    assert placement.previous_text == "Traffic enters through EDGE-01."
    assert placement.following_text == "Controllers run in separate availability zones."


def test_pptx_picture_extractor_keeps_each_slide_occurrence(tmp_path):
    from src.ingestion.figure_extract import extract_image_regions_pptx

    path = tmp_path / "repeated.pptx"
    _make_pptx(path, repeated_slide=True)
    regions = extract_image_regions_pptx(path)

    assert [region.figure_id for region in regions] == ["s1-fig-001", "s2-fig-001"]
    assert [region.page for region in regions] == [0, 1]
    assert regions[0].content_hash == regions[1].content_hash
    assert regions[0].source == "pptx"
    assert regions[0].image_bytes.startswith(b"\x89PNG\r\n\x1a\n")
    assert regions[0].caption == "Figure 2. Controller topology"


def test_pptx_enrichment_inserts_description_at_picture_anchor(tmp_path, monkeypatch):
    from src.config import settings
    from src.ingestion.figure_extract import (
        FigureEnrichmentResult,
        FigureRecord,
        ProseBlock,
        enrich_office_document_with_figures,
    )

    path = tmp_path / "architecture.pptx"
    _make_pptx(path)
    monkeypatch.setattr(settings, "figure_extraction_enabled", True)

    def fake_process(regions, **kwargs):
        region = regions[0]
        description = (
            f"[Figure {region.figure_id} — network]\n"
            "Nodes: CTRL-01, CTRL-02\n[/Figure]"
        )
        return FigureEnrichmentResult(
            prose_blocks=[ProseBlock(description, page=region.page)],
            figure_records=[FigureRecord(
                figure_id=region.figure_id,
                description=description,
                kind="network",
                content_hash=region.content_hash,
                body_index=region.body_index,
                section_path=region.section_path,
                caption=region.caption,
                alt_text=region.alt_text,
                previous_text=region.previous_text,
                following_text=region.following_text,
                source="pptx",
                slide=region.page,
            )],
            figures_seen=1,
            figures_used=1,
        )

    with patch("src.ingestion.figure_extract.process_image_regions", side_effect=fake_process):
        enriched = enrich_office_document_with_figures(path, "unused flat text")

    assert "## Embedded figures" not in enriched.enriched_text
    assert enriched.enriched_text.index("EDGE-01") < enriched.enriched_text.index("CTRL-01")
    assert enriched.enriched_text.index("CTRL-02") < enriched.enriched_text.index("Figure 2")
    assert enriched.figures[0].slide == 0


def test_pptx_parser_extracts_table_and_chart_data(tmp_path):
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
    from src.ingestion.parser import parse_document

    path = tmp_path / "data.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Capacity"
    table = slide.shapes.add_table(2, 2, Inches(0.5), Inches(1), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Nodes"
    table.cell(1, 0).text = "East"
    table.cell(1, 1).text = "12"
    data = ChartData()
    data.categories = ["East", "West"]
    data.add_series("Controllers", (12, 9))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(5), Inches(1), Inches(4), Inches(3), data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Controller Count"
    presentation.save(path)

    parsed = parse_document(path)
    assert "Region | Nodes" in parsed.text
    assert "East | 12" in parsed.text
    assert "Title: Controller Count" in parsed.text
    assert "Categories: East | West" in parsed.text
    assert "Series Controllers: 12.0 | 9.0" in parsed.text
