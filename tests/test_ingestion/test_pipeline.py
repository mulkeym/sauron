import pytest
from pathlib import Path
from unittest.mock import MagicMock, AsyncMock, patch, call

from src.ingestion.pipeline import IngestResult, ingest_document
from src.knowledge.categorizer import CategorizationResult

FIXTURES = Path(__file__).parent.parent.parent / "test_fixtures"


@pytest.fixture
def mock_deps():
    mock_vector_store = MagicMock()
    mock_metadata_store = AsyncMock()
    mock_embed = MagicMock(return_value=[[0.1] * 1024])
    return mock_vector_store, mock_metadata_store, mock_embed


@pytest.mark.asyncio
async def test_ingest_pdf(mock_deps):
    vector_store, metadata_store, mock_embed = mock_deps
    with patch("src.ingestion.pipeline.embed_texts", mock_embed):
        result = await ingest_document(
            file_path=FIXTURES / "sample.pdf",
            acl_groups=["finance"],
            uploaded_by="mike",
            vector_store=vector_store,
            metadata_store=metadata_store,
        )
    assert isinstance(result, IngestResult)
    assert result.doc_type == "pdf"
    assert result.chunk_count > 0
    vector_store.upsert.assert_called_once()
    metadata_store.add_document.assert_called_once()


@pytest.mark.asyncio
async def test_ingest_transcript(mock_deps):
    vector_store, metadata_store, mock_embed = mock_deps
    with patch("src.ingestion.pipeline.embed_texts", mock_embed):
        result = await ingest_document(
            file_path=FIXTURES / "sample_transcript.txt",
            acl_groups=["engineering"],
            uploaded_by="mike",
            vector_store=vector_store,
            metadata_store=metadata_store,
        )
    assert result.doc_type == "transcript"
    assert result.chunk_count > 0


@pytest.mark.asyncio
async def test_ingest_auto_categorizes(mock_deps):
    vector_store, metadata_store, mock_embed = mock_deps
    metadata_store.list_categories = AsyncMock(return_value=[])
    metadata_store.add_proposal = AsyncMock()

    mock_cat_result = CategorizationResult(category="finance_policies", confidence=0.9, is_new=False)

    with patch("src.ingestion.pipeline.embed_texts", mock_embed):
        with patch("src.ingestion.pipeline.categorize_document", return_value=mock_cat_result):
            result = await ingest_document(
                file_path=FIXTURES / "sample.pdf", acl_groups=["finance"],
                uploaded_by="mike", vector_store=vector_store, metadata_store=metadata_store,
                auto_categorize=True,
            )
    assert result.doc_type == "pdf"
    metadata_store.add_document.assert_called_once()


@pytest.mark.asyncio
async def test_ingest_extracts_entities(mock_deps):
    vector_store, metadata_store, mock_embed = mock_deps
    metadata_store.add_entity = AsyncMock(return_value=1)
    metadata_store.add_mention = AsyncMock()
    metadata_store.add_relationship = AsyncMock()

    from src.knowledge.extractor import ExtractionResult
    mock_extraction = ExtractionResult(
        entities=[{"name": "Mike", "type": "person"}, {"name": "Policy 4.2", "type": "policy"}],
        relationships=[{"source": "Policy 4.2", "target": "expense reporting", "type": "governs"}],
        sections=[],
    )

    with patch("src.ingestion.pipeline.embed_texts", mock_embed):
        with patch("src.ingestion.pipeline.extract_entities", return_value=mock_extraction):
            result = await ingest_document(
                file_path=FIXTURES / "sample.pdf", acl_groups=["finance"],
                uploaded_by="mike", vector_store=vector_store, metadata_store=metadata_store,
            )
    assert result.chunk_count > 0
    metadata_store.add_entity.assert_called()
    metadata_store.add_mention.assert_called()


@pytest.mark.asyncio
async def test_ingest_document_spreadsheet_dedups_clean_sheet(tmp_path, monkeypatch):
    """A clean spreadsheet that is structured-ingested emits NO full-text tier
    chunks (only its DuckDB rows + table_row narratives); messy sheets would
    still produce structure-aware text chunks."""
    import openpyxl
    from src.ingestion import pipeline

    p = tmp_path / "pay.xlsx"
    wb = openpyxl.Workbook(); wb.remove(wb.active)
    clean = wb.create_sheet("Pay")
    for r in [["locality", "grade", "salary"], ["Tampa", "GS-12", 86415],
              ["Boston", "GS-12", 92000], ["Denver", "GS-13", 99000]]:
        clean.append(r)
    wb.save(p)

    captured = {"tiers": []}

    class FakeVS:
        def upsert(self, texts, vectors, metadatas):
            captured["tiers"] += [m.chunk_size_tier for m in metadatas]

    class FakeMS:
        async def add_document(self, **k): pass
        async def get_category(self, name): return None
        async def add_category(self, **k): pass
        async def save_schema(self, s): pass

    # Stub everything external: embeddings (both call sites), the LLM summary/
    # profiler, the LightRAG insert, and the schema registry.
    monkeypatch.setattr(pipeline, "embed_texts", lambda texts, *a, **k: [[0.0] for _ in texts])
    monkeypatch.setattr("src.ingestion.tabular_ingest.embed_texts",
                        lambda texts, *a, **k: [[0.0] for _ in texts])
    monkeypatch.setattr("src.generation.llm_client.generate", lambda **k: "")

    async def _coro(*a, **k):
        return None
    monkeypatch.setattr("src.knowledge.graph_rag.insert_document", lambda *a, **k: _coro())
    monkeypatch.setattr("src.api.routes_ingest.get_schema_registry",
                        lambda: type("R", (), {"register": lambda self, s: None,
                                               "remove": lambda self, *a: None})())

    await pipeline.ingest_document(
        str(p), acl_groups=["g1"], uploaded_by="t", vector_store=FakeVS(),
        metadata_store=FakeMS(), category="cat",
    )
    # The clean "Pay" sheet was de-duped: no small/medium/large/xlarge text tiers,
    # only table_row narratives.
    assert "table_row" in captured["tiers"]
    assert not ({"small", "medium", "large", "xlarge"} & set(captured["tiers"]))


def _stub_pipeline(monkeypatch, pipeline, kg_calls):
    """Stub the external deps shared by the KG-gating tests, recording every
    LightRAG insert_document call into kg_calls."""
    async def _coro(*a, **k):
        return None

    def _fake_insert(*a, **k):
        kg_calls.append((a, k))
        return _coro()

    monkeypatch.setattr(pipeline, "embed_texts", lambda texts, *a, **k: [[0.0] for _ in texts])
    monkeypatch.setattr("src.ingestion.tabular_ingest.embed_texts",
                        lambda texts, *a, **k: [[0.0] for _ in texts])
    monkeypatch.setattr("src.generation.llm_client.generate", lambda **k: "")
    monkeypatch.setattr("src.knowledge.graph_rag.insert_document", _fake_insert)
    monkeypatch.setattr("src.api.routes_ingest.get_schema_registry",
                        lambda: type("R", (), {"register": lambda self, s: None,
                                               "remove": lambda self, *a: None})())


class _GateVS:
    def upsert(self, texts, vectors, metadatas): pass


class _GateMS:
    async def add_document(self, **k): pass
    async def get_category(self, name): return None
    async def add_category(self, **k): pass
    async def save_schema(self, s): pass


@pytest.mark.asyncio
async def test_ingest_spreadsheet_skips_knowledge_graph(tmp_path, monkeypatch):
    """Spreadsheets are fully covered by the structured/tabular store, so KG
    extraction (costly + near-empty for numeric tables) is skipped."""
    import openpyxl
    from src.ingestion import pipeline

    p = tmp_path / "pay.xlsx"
    wb = openpyxl.Workbook(); wb.remove(wb.active)
    s = wb.create_sheet("Pay")
    for r in [["locality", "grade", "salary"], ["Tampa", "GS-12", 86415]]:
        s.append(r)
    wb.save(p)

    kg_calls = []
    _stub_pipeline(monkeypatch, pipeline, kg_calls)
    await pipeline.ingest_document(str(p), acl_groups=["g1"], uploaded_by="t",
                                   vector_store=_GateVS(), metadata_store=_GateMS(), category="cat")
    assert kg_calls == []   # KG skipped for spreadsheet


@pytest.mark.asyncio
async def test_ingest_plaintext_still_builds_knowledge_graph(tmp_path, monkeypatch):
    """Regression guard: the spreadsheet gate must NOT suppress KG for prose docs."""
    from src.ingestion import pipeline

    p = tmp_path / "memo.txt"
    p.write_text("Acme Corp signed a contract with Globex in 2026.")

    kg_calls = []
    _stub_pipeline(monkeypatch, pipeline, kg_calls)
    await pipeline.ingest_document(str(p), acl_groups=["g1"], uploaded_by="t",
                                   vector_store=_GateVS(), metadata_store=_GateMS(), category="cat")
    assert len(kg_calls) == 1   # KG still runs for non-spreadsheet


@pytest.mark.asyncio
async def test_docx_figure_is_inline_in_kg_and_has_dedicated_chunk(tmp_path, monkeypatch):
    from docx import Document
    from docx.shared import Inches
    from PIL import Image
    from src.ingestion import pipeline
    from src.ingestion.figure_extract import (
        FigureEnrichmentResult, FigureRecord, ProseBlock,
    )

    image_path = tmp_path / "topology.png"
    Image.new("RGB", (600, 400), "navy").save(image_path)
    docx_path = tmp_path / "manual.docx"
    doc = Document()
    doc.add_heading("Control Plane", level=1)
    doc.add_paragraph("Traffic enters through BRANCH-01.")
    doc.add_picture(str(image_path), width=Inches(4))
    doc.add_paragraph("Controllers run in separate zones.")
    doc.save(docx_path)

    def _fake_process(regions, **kwargs):
        region = regions[0]
        description = (
            f"[Figure {region.figure_id} — network]\n"
            "Nodes:\n- CTRL-01\n- CTRL-02\n"
            "Links:\n- CTRL-01 --sync--> CTRL-02\n[/Figure]"
        )
        return FigureEnrichmentResult(
            prose_blocks=[ProseBlock(text=description, page=0)],
            figure_records=[FigureRecord(
                figure_id=region.figure_id, description=description, kind="network",
                content_hash=region.content_hash, body_index=region.body_index,
                section_path=region.section_path,
                previous_text=region.previous_text, following_text=region.following_text,
            )],
            figures_seen=1, figures_used=1,
        )

    captured = []

    class FigureVS:
        def upsert(self, texts, vectors, metadatas):
            captured.extend(zip(texts, metadatas))

    kg_calls = []
    _stub_pipeline(monkeypatch, pipeline, kg_calls)
    monkeypatch.setattr("src.ingestion.figure_extract.process_image_regions", _fake_process)

    result = await pipeline.ingest_document(
        str(docx_path), acl_groups=["g1"], uploaded_by="t",
        vector_store=FigureVS(), metadata_store=_GateMS(), category="cat",
    )

    figure_chunks = [(text, meta) for text, meta in captured if meta.content_type == "figure"]
    assert len(figure_chunks) == 1
    figure_text, figure_meta = figure_chunks[0]
    assert figure_meta.figure_id == "fig-001"
    assert figure_meta.section_title == "Control Plane"
    assert "Context before: Traffic enters through BRANCH-01." in figure_text
    assert result.chunk_count > 1

    assert len(kg_calls) == 1
    kg_text = kg_calls[0][0][0]
    assert kg_text.index("BRANCH-01") < kg_text.index("CTRL-01")
    assert kg_text.index("CTRL-02") < kg_text.index("Controllers run")


@pytest.mark.asyncio
async def test_pdf_figure_has_page_chunk_and_ordered_kg(monkeypatch):
    from src.ingestion import pipeline
    from src.ingestion.figure_extract import FigureRecord
    from src.ingestion.pdf_extract import ExtractedPdf, ProseBlock

    record = FigureRecord(
        figure_id="p1-fig-001",
        description="[Figure p1-fig-001 — network]\n- EDGE-01 --> CORE-01\n[/Figure]",
        kind="network", page=0, bbox=(80, 120, 500, 300),
        source="embedded", section_path=["Network Architecture"],
        previous_text="Traffic enters through EDGE-01.",
        following_text="CORE-01 forwards traffic to the service tier.",
    )
    base = ExtractedPdf(
        prose_blocks=[ProseBlock("flat page text", page=0)], method="digital",
    )
    enriched = ExtractedPdf(
        prose_blocks=[
            ProseBlock("Traffic enters through EDGE-01.", page=0),
            ProseBlock(record.description, page=0, content_type="figure",
                       figure_id=record.figure_id),
            ProseBlock("CORE-01 forwards traffic to the service tier.", page=0),
        ],
        method="mixed", figure_records=[record],
    )
    monkeypatch.setattr(pipeline, "extract_pdf", lambda path: base)
    monkeypatch.setattr(
        "src.ingestion.figure_extract.enrich_pdf_with_figures",
        lambda path, extracted: enriched,
    )

    captured = []

    class FigureVS:
        def upsert(self, texts, vectors, metadatas):
            captured.extend(zip(texts, metadatas))

    kg_calls = []
    _stub_pipeline(monkeypatch, pipeline, kg_calls)
    result = await pipeline.ingest_document(
        str(FIXTURES / "sample.pdf"), acl_groups=["g1"], uploaded_by="t",
        vector_store=FigureVS(), metadata_store=_GateMS(), category="cat",
    )

    figure_chunks = [(text, meta) for text, meta in captured if meta.content_type == "figure"]
    assert len(figure_chunks) == 1
    _, meta = figure_chunks[0]
    assert meta.figure_id == "p1-fig-001"
    assert meta.page == 1
    assert meta.section_title == "Network Architecture"
    assert "page 1" in meta.source_locator
    assert result.chunk_count > 1

    kg_text = kg_calls[0][0][0]
    assert kg_text.index("Traffic enters") < kg_text.index("EDGE-01 --> CORE-01")
    assert kg_text.index("EDGE-01 --> CORE-01") < kg_text.index("forwards traffic")


@pytest.mark.asyncio
async def test_pptx_figure_has_slide_chunk_and_ordered_kg(tmp_path, monkeypatch):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches
    from src.ingestion import pipeline
    from src.ingestion.figure_extract import (
        FigureEnrichmentResult, FigureRecord, ProseBlock,
    )

    image_path = tmp_path / "topology.png"
    Image.new("RGB", (600, 400), "navy").save(image_path)
    pptx_path = tmp_path / "manual.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Control Plane"
    before = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(6), Inches(0.5))
    before.text = "Traffic enters through EDGE-01."
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(4))
    after = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(7), Inches(0.5))
    after.text = "CORE-01 forwards traffic to the service tier."
    presentation.save(pptx_path)

    def fake_process(regions, **kwargs):
        region = regions[0]
        description = (
            f"[Figure {region.figure_id} — network]\n"
            "- EDGE-01 --> CORE-01\n[/Figure]"
        )
        return FigureEnrichmentResult(
            prose_blocks=[ProseBlock(description, page=region.page)],
            figure_records=[FigureRecord(
                figure_id=region.figure_id, description=description, kind="network",
                content_hash=region.content_hash, body_index=region.body_index,
                section_path=region.section_path, previous_text=region.previous_text,
                following_text=region.following_text, source="pptx", slide=region.page,
            )],
            figures_seen=1, figures_used=1,
        )

    captured = []

    class FigureVS:
        def upsert(self, texts, vectors, metadatas):
            captured.extend(zip(texts, metadatas))

    kg_calls = []
    _stub_pipeline(monkeypatch, pipeline, kg_calls)
    monkeypatch.setattr("src.ingestion.figure_extract.process_image_regions", fake_process)

    result = await pipeline.ingest_document(
        str(pptx_path), acl_groups=["g1"], uploaded_by="t",
        vector_store=FigureVS(), metadata_store=_GateMS(), category="cat",
    )

    figure_chunks = [(text, meta) for text, meta in captured if meta.content_type == "figure"]
    assert len(figure_chunks) == 1
    figure_text, meta = figure_chunks[0]
    assert meta.figure_id == "s1-fig-001"
    assert meta.slide == 1
    assert meta.section_title == "Control Plane"
    assert "slide 1" in meta.source_locator
    assert "Context before: Traffic enters through EDGE-01." in figure_text
    assert result.chunk_count > 1

    kg_text = kg_calls[0][0][0]
    assert kg_text.index("Traffic enters") < kg_text.index("EDGE-01 --> CORE-01")
    assert kg_text.index("EDGE-01 --> CORE-01") < kg_text.index("forwards traffic")
