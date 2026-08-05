"""Extract embedded / page images from PDFs and turn them into text or tables.

Pipeline:
  1. Collect image regions (embedded PdfImage objects + sparse-page renders)
  2. Skip logos / tiny images
  3. OCR (optional, first)
  4. Classify: table | network | process | text_scan | other
  5. Strategy → SheetGrid and/or ProseBlock
  6. Merge into ExtractedPdf before chunking / DuckDB ingest

Fail-open: any vision/OCR failure logs and continues so ingest never dies on one figure.
"""
from __future__ import annotations

import hashlib
import io
import logging
import re
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path

from src.config import settings
from src.ingestion.pdf_extract import ExtractedPdf, ProseBlock, normalize_grid
from src.ingestion.tabular import SheetGrid

logger = logging.getLogger(__name__)


class ImageKind(str, Enum):
    TABLE = "table"
    NETWORK = "network"
    PROCESS = "process"
    TEXT_SCAN = "text_scan"
    OTHER = "other"
    SKIP = "skip"


@dataclass
class ImageRegion:
    page: int
    index: int
    image_bytes: bytes
    width: int
    height: int
    source: str = "embedded"  # "embedded" | "page_render"
    ocr_text: str = ""
    kind: ImageKind = ImageKind.OTHER
    content_hash: str = ""
    figure_id: str = ""
    relationship_id: str = ""
    body_index: int | None = None
    section_path: list[str] = field(default_factory=list)
    caption: str = ""
    alt_text: str = ""
    previous_text: str = ""
    following_text: str = ""
    bbox: tuple[float, float, float, float] | None = None


@dataclass
class FigureRecord:
    """Searchable analysis and provenance for one figure occurrence."""

    figure_id: str
    description: str
    kind: str
    content_hash: str = ""
    body_index: int | None = None
    section_path: list[str] = field(default_factory=list)
    caption: str = ""
    alt_text: str = ""
    previous_text: str = ""
    following_text: str = ""
    page: int | None = None
    bbox: tuple[float, float, float, float] | None = None
    source: str = ""
    slide: int | None = None

    def retrieval_text(self) -> str:
        parts = [f"Figure: {self.figure_id}"]
        if self.page is not None:
            parts.append(f"Page: {self.page + 1}")
        if self.slide is not None:
            parts.append(f"Slide: {self.slide + 1}")
        if self.section_path:
            parts.append("Section: " + " > ".join(self.section_path))
        if self.caption:
            parts.append(f"Caption: {self.caption}")
        if self.alt_text:
            parts.append(f"Alt text: {self.alt_text}")
        if self.previous_text:
            parts.append(f"Context before: {self.previous_text}")
        parts.append(self.description.strip())
        if self.following_text:
            parts.append(f"Context after: {self.following_text}")
        return "\n".join(p for p in parts if p.strip())


@dataclass
class OfficeFigureResult:
    enriched_text: str
    table_grids: list[SheetGrid] = field(default_factory=list)
    figures: list[FigureRecord] = field(default_factory=list)


@dataclass
class FigureEnrichmentResult:
    """Extras produced by figure extraction (merged into ExtractedPdf by caller)."""
    table_grids: list[SheetGrid] = field(default_factory=list)
    prose_blocks: list[ProseBlock] = field(default_factory=list)
    figures_seen: int = 0
    figures_used: int = 0
    figures_skipped: int = 0
    figure_records: list[FigureRecord] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Image collection
# ---------------------------------------------------------------------------

def _pil_to_png_bytes(img) -> bytes:
    buf = io.BytesIO()
    # Normalize mode for JPEG-sourced images etc.
    if img.mode not in ("RGB", "L"):
        img = img.convert("RGB")
    img.save(buf, format="PNG")
    return buf.getvalue()


def _should_skip_size(width: int, height: int) -> bool:
    if width < settings.figure_min_width or height < settings.figure_min_height:
        return True
    if width * height < settings.figure_min_area:
        return True
    # Extreme aspect ratio (page rules / lines)
    ratio = max(width, height) / max(1, min(width, height))
    if ratio > 12:
        return True
    # Small square icons (UI chrome) — e.g. 128–160px pencils — waste vision budget
    if width <= 160 and height <= 160 and abs(width - height) <= 24:
        return True
    return False


def _region_from_image_bytes(
    raw: bytes,
    *,
    page: int,
    index: int,
    source: str,
    seen_hashes: set[str],
) -> ImageRegion | None:
    """Decode arbitrary image bytes to a PNG ImageRegion, or None if skip/fail."""
    try:
        from PIL import Image

        img = Image.open(io.BytesIO(raw))
        img.load()
        if _should_skip_size(img.width, img.height):
            return None
        png = _pil_to_png_bytes(img)
        digest = hashlib.sha256(png).hexdigest()[:16]
        if digest in seen_hashes:
            return None
        seen_hashes.add(digest)
        return ImageRegion(
            page=page,
            index=index,
            image_bytes=png,
            width=img.width,
            height=img.height,
            source=source,
            content_hash=digest,
        )
    except Exception as e:
        logger.debug(f"figure extract: decode image failed ({source}): {e}")
        return None


def extract_image_regions_from_zip_media(
    path: Path,
    *,
    media_prefixes: tuple[str, ...],
    source: str,
) -> list[ImageRegion]:
    """Pull images from OOXML packages (docx/xlsx/pptx) via zip media paths."""
    import zipfile

    path = Path(path)
    regions: list[ImageRegion] = []
    seen: set[str] = set()
    if not zipfile.is_zipfile(path):
        return regions
    try:
        with zipfile.ZipFile(path) as zf:
            names = [
                n for n in zf.namelist()
                if any(n.startswith(p) for p in media_prefixes)
                and not n.endswith("/")
            ]
            names.sort()
            for i, name in enumerate(names):
                try:
                    raw = zf.read(name)
                except Exception:
                    continue
                # sheet/page index unknown from zip alone — use sequential order
                reg = _region_from_image_bytes(
                    raw, page=0, index=i, source=source, seen_hashes=seen,
                )
                if reg:
                    regions.append(reg)
    except Exception as e:
        logger.warning(f"figure extract: zip media scan failed for {path.name}: {e}")
    return regions


def extract_image_regions_docx(path: Path) -> list[ImageRegion]:
    """Embedded Word images with their actual body placements and context."""
    path = Path(path)
    # Resolve relationship IDs captured by the ordered Word parser.  A fresh
    # decode set per placement deliberately retains repeated occurrences.
    try:
        from docx import Document
        from src.ingestion.parser import parse_document

        doc = Document(str(path))
        ordered: list[ImageRegion] = []
        parsed = parse_document(path)
        placements = [
            b.figure for b in parsed.blocks
            if b.block_type == "figure" and b.figure is not None
        ]
        for idx, placement in enumerate(placements):
            rel = doc.part.rels.get(placement.relationship_id)
            if rel is None or "image" not in (rel.reltype or ""):
                continue
            try:
                raw = rel.target_part.blob
            except Exception:
                continue
            reg = _region_from_image_bytes(
                raw, page=0, index=idx, source="docx", seen_hashes=set(),
            )
            if reg:
                reg.figure_id = placement.figure_id
                reg.relationship_id = placement.relationship_id
                reg.body_index = placement.body_index
                reg.section_path = list(placement.section_path)
                reg.caption = placement.caption
                reg.alt_text = placement.alt_text
                reg.previous_text = placement.previous_text
                reg.following_text = placement.following_text
                ordered.append(reg)
        if ordered:
            return ordered
    except Exception as e:
        logger.debug(f"figure extract: ordered DOCX walk failed: {e}")

    regions = extract_image_regions_from_zip_media(
        path, media_prefixes=("word/media/",), source="docx",
    )
    # Compatibility fallback for malformed documents without body anchors.
    for i, region in enumerate(regions):
        region.figure_id = f"fig-{i + 1:03d}"
    return regions


def extract_image_regions_xlsx(path: Path) -> list[ImageRegion]:
    """Embedded images from Excel .xlsx/.xlsm (sheet drawings + xl/media)."""
    path = Path(path)
    regions: list[ImageRegion] = []
    seen: set[str] = set()
    # 1) openpyxl worksheet drawings (preserves sheet index as "page")
    try:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), data_only=False)
        for sheet_i, ws in enumerate(wb.worksheets):
            images = list(getattr(ws, "_images", None) or [])
            for img_i, img in enumerate(images):
                try:
                    raw = img._data() if hasattr(img, "_data") else None
                    if not raw and getattr(img, "ref", None):
                        # path-like ref
                        from PIL import Image as PILImage
                        pil = PILImage.open(img.ref)
                        raw = _pil_to_png_bytes(pil)
                    if not raw:
                        continue
                    reg = _region_from_image_bytes(
                        raw if isinstance(raw, (bytes, bytearray)) else bytes(raw),
                        page=sheet_i,
                        index=img_i,
                        source="xlsx",
                        seen_hashes=seen,
                    )
                    if reg:
                        regions.append(reg)
                except Exception as e:
                    logger.debug(f"figure extract: xlsx image on sheet {sheet_i}: {e}")
    except Exception as e:
        logger.debug(f"figure extract: openpyxl image walk failed: {e}")

    # 2) zip media fallback (catches images not linked via _images)
    zip_regions = extract_image_regions_from_zip_media(
        path, media_prefixes=("xl/media/",), source="xlsx",
    )
    for reg in zip_regions:
        if reg.content_hash not in seen:
            # re-add with next index on page 0 if new
            seen.add(reg.content_hash)
            reg.index = len(regions)
            regions.append(reg)
    return regions


def extract_image_regions_pptx(path: Path) -> list[ImageRegion]:
    """PowerPoint pictures resolved from slide-local relationship IDs."""
    path = Path(path)
    try:
        from pptx import Presentation
        from src.ingestion.parser import parse_document

        presentation = Presentation(str(path))
        parsed = parse_document(path)
        ordered: list[ImageRegion] = []
        placements = [
            block.figure for block in parsed.blocks
            if block.block_type == "figure" and block.figure is not None
        ]
        for placement in placements:
            if placement.page is None or placement.page >= len(presentation.slides):
                continue
            slide = presentation.slides[placement.page]
            rel = slide.part.rels.get(placement.relationship_id)
            if rel is None or "image" not in (rel.reltype or ""):
                continue
            try:
                raw = rel.target_part.blob
            except Exception:
                continue
            reg = _region_from_image_bytes(
                raw, page=placement.page, index=len(ordered), source="pptx",
                seen_hashes=set(),
            )
            if reg:
                reg.figure_id = placement.figure_id
                reg.relationship_id = placement.relationship_id
                reg.body_index = placement.body_index
                reg.section_path = list(placement.section_path)
                reg.caption = placement.caption
                reg.alt_text = placement.alt_text
                reg.previous_text = placement.previous_text
                reg.following_text = placement.following_text
                reg.bbox = placement.bbox
                ordered.append(reg)
        if ordered:
            return ordered
    except Exception as e:
        logger.debug(f"figure extract: ordered PPTX walk failed: {e}")

    regions = extract_image_regions_from_zip_media(
        path, media_prefixes=("ppt/media/",), source="pptx",
    )
    for i, region in enumerate(regions):
        region.figure_id = f"s1-fig-{i + 1:03d}"
    return regions


def extract_image_regions(path: Path) -> list[ImageRegion]:
    """Collect embedded images and optional full-page renders from a PDF."""
    import pypdfium2 as pdfium
    from pypdfium2.raw import FPDF_PAGEOBJ_IMAGE

    path = Path(path)
    regions: list[ImageRegion] = []
    seen_hashes: set[str] = set()

    try:
        pdf = pdfium.PdfDocument(str(path))
    except Exception as e:
        logger.warning(f"figure extract: cannot open PDF {path.name}: {e}")
        return regions

    try:
        for page_idx in range(len(pdf)):
            page = pdf[page_idx]
            img_idx = 0
            try:
                page_width, page_height = page.get_size()
            except Exception:
                page_width = page_height = 0

            # --- Embedded image objects ---
            try:
                for obj in page.get_objects():
                    try:
                        if getattr(obj, "type", None) != FPDF_PAGEOBJ_IMAGE:
                            continue
                        # PdfImage
                        try:
                            w, h = obj.get_px_size()
                        except Exception:
                            w = h = 0
                        if _should_skip_size(int(w or 0), int(h or 0)):
                            continue
                        try:
                            bitmap = obj.get_bitmap(render=True)
                            pil = bitmap.to_pil()
                            bitmap.close()
                        except Exception as e:
                            logger.debug(f"figure extract: bitmap failed p{page_idx}: {e}")
                            continue
                        if _should_skip_size(pil.width, pil.height):
                            continue
                        png = _pil_to_png_bytes(pil)
                        digest = hashlib.sha256(png).hexdigest()[:16]
                        bbox = None
                        try:
                            left, bottom, right, top = obj.get_bounds()
                            # PDFium uses a bottom-left origin; pdfplumber's text
                            # layout uses a top-left origin.
                            bbox = (
                                float(min(left, right)),
                                float(page_height - max(top, bottom)),
                                float(max(left, right)),
                                float(page_height - min(top, bottom)),
                            )
                        except Exception:
                            pass
                        regions.append(
                            ImageRegion(
                                page=page_idx,
                                index=img_idx,
                                image_bytes=png,
                                width=pil.width,
                                height=pil.height,
                                source="embedded",
                                content_hash=digest,
                                bbox=bbox,
                            )
                        )
                        img_idx += 1
                    except Exception as e:
                        logger.debug(f"figure extract: skip object on p{page_idx}: {e}")
            except Exception as e:
                logger.debug(f"figure extract: get_objects failed p{page_idx}: {e}")

            # --- Sparse digital text → full-page render (diagram-as-page) ---
            if settings.figure_render_text_sparse_pages:
                try:
                    textpage = page.get_textpage()
                    page_text = (textpage.get_text_bounded() or "").strip()
                    textpage.close()
                except Exception:
                    page_text = ""
                if len(page_text) < settings.figure_sparse_text_chars:
                    # Avoid double-counting if we already have a large embedded image
                    has_large = any(
                        r.page == page_idx and r.width * r.height > 200_000
                        for r in regions
                    )
                    if not has_large:
                        try:
                            scale = settings.figure_page_render_dpi_scale
                            pil = page.render(scale=scale).to_pil()
                            if not _should_skip_size(pil.width, pil.height):
                                png = _pil_to_png_bytes(pil)
                                digest = hashlib.sha256(png).hexdigest()[:16]
                                if digest not in seen_hashes:
                                    seen_hashes.add(digest)
                                    regions.append(
                                        ImageRegion(
                                            page=page_idx,
                                            index=img_idx,
                                            image_bytes=png,
                                            width=pil.width,
                                            height=pil.height,
                                            source="page_render",
                                            content_hash=digest,
                                            bbox=(
                                                0.0, 0.0,
                                                float(page_width or pil.width),
                                                float(page_height or pil.height),
                                            ),
                                        )
                                    )
                        except Exception as e:
                            logger.debug(f"figure extract: page render failed p{page_idx}: {e}")
    finally:
        try:
            pdf.close()
        except Exception:
            pass

    regions.sort(key=lambda r: (
        r.page,
        r.bbox[1] if r.bbox else float("inf"),
        r.bbox[0] if r.bbox else r.index,
    ))
    per_page: dict[int, int] = {}
    for body_index, region in enumerate(regions):
        per_page[region.page] = per_page.get(region.page, 0) + 1
        region.index = per_page[region.page] - 1
        region.figure_id = f"p{region.page + 1}-fig-{per_page[region.page]:03d}"
        region.body_index = body_index
    return regions


def extract_image_regions_for_path(path: Path) -> list[ImageRegion]:
    """Dispatch image collection by file extension."""
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_image_regions(path)
    if suffix == ".docx":
        return extract_image_regions_docx(path)
    if suffix in (".xlsx", ".xlsm"):
        return extract_image_regions_xlsx(path)
    if suffix == ".pptx":
        return extract_image_regions_pptx(path)
    logger.info(f"figure extract: no image collector for {suffix}")
    return []


# ---------------------------------------------------------------------------
# OCR + classification
# ---------------------------------------------------------------------------

_IP_RE = re.compile(r"\b(?:\d{1,3}\.){3}\d{1,3}\b")
_HOST_RE = re.compile(r"\b[a-zA-Z0-9][\w.-]*\.(?:local|corp|com|net|io|internal)\b", re.I)
_URL_RE = re.compile(r"https?://[^\s]+", re.I)
# Tokens worth preserving from OCR when vision omits them
_OCR_TOKEN_RE = re.compile(
    r"(?:"
    r"https?://[^\s]+|"
    r"\b(?:\d{1,3}\.){3}\d{1,3}(?:/\d{1,2})?\b|"
    r"\b[A-Z][A-Za-z0-9]+(?:[ -][A-Z][A-Za-z0-9]+){1,5}\b|"  # Title Case phrases
    r"\b[A-Z]{2,}[A-Z0-9]*\b|"  # ACRONYMS RBAC TACACS VRF
    r"\b(?:vManage|vBond|vSmart|vMonitor|vOrchestrator|vdaemon|SD-WAN|IPsec|DTLS|SAML|Okta|Jenkins|Bamboo)\b"
    r")",
    re.I,
)

_DIAGRAM_HINTS = (
    "portal", "overlay", "architecture", "lifecycle", "topology", "workflow",
    "pipeline", "controller", "deploy", "accelerator", "regulation", "visibility",
    "operation", "services", "network", "router", "cloud",
)
_PROCESS_HINTS = (
    "workflow", "pipeline", "lifecycle", "process", "portal", "phase",
    "jenkins", "bamboo", "prepare", "incident", "decision", "step",
)
_NETWORK_HINTS = (
    "router", "overlay", "tloc", "ipsec", "dtls", "vlan", "vrf", "topology",
    "firewall", "controller", "validator", "vmanage", "vbond",
)


def ocr_image(image_bytes: bytes) -> str:
    """Run tesseract OCR; return empty string on any failure."""
    try:
        import pytesseract
        from PIL import Image

        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img) or ""
        return text.strip()
    except Exception as e:
        logger.debug(f"OCR failed: {e}")
        return ""


def _ocr_line_stats(text: str) -> dict:
    lines = [ln.strip() for ln in (text or "").splitlines() if ln.strip()]
    n = len(lines)
    if not n:
        return {
            "lines": [], "n_lines": 0, "avg_len": 0.0,
            "short_lines": 0, "long_lines": 0, "chars": 0,
        }
    lengths = [len(ln) for ln in lines]
    return {
        "lines": lines,
        "n_lines": n,
        "avg_len": sum(lengths) / n,
        "short_lines": sum(1 for L in lengths if L < 40),
        "long_lines": sum(1 for L in lengths if L >= 55),
        "chars": sum(lengths),
    }


def looks_like_linear_prose(text: str) -> bool:
    """True when OCR reads as flowing paragraphs (OCR-only is enough)."""
    st = _ocr_line_stats(text)
    if st["n_lines"] < 2 or st["chars"] < 80:
        return False
    # Mostly long lines, high average length → document screenshot / prose
    if st["avg_len"] >= 48 and st["long_lines"] >= max(2, st["n_lines"] * 0.45):
        return True
    # Few lines but each is a long sentence
    if st["n_lines"] <= 4 and st["avg_len"] >= 70:
        return True
    return False


def looks_like_diagram(region: ImageRegion, text: str) -> bool:
    """Spatial / multi-panel figure — needs vision even if OCR found many words."""
    st = _ocr_line_stats(text)
    w, h = region.width, region.height
    large = w >= 550 and h >= 220
    wide = w >= 700 and h >= 180
    # Many short label lines = layout diagram (p.22 portal graphic)
    if (large or wide) and st["n_lines"] >= 4 and st["short_lines"] >= 3 and st["avg_len"] < 50:
        return True
    if large and st["n_lines"] >= 6 and st["avg_len"] < 55:
        return True
    low = (text or "").lower()
    hint_hits = sum(1 for k in _DIAGRAM_HINTS if k in low)
    if (large or wide) and hint_hits >= 2 and st["avg_len"] < 60:
        return True
    # Sparse OCR on a large image alone is NOT enough (often photos/art).
    # Require layout/keyword cues so pure illustrations use the illustration path.
    if large and st["chars"] < 40 and hint_hits >= 1:
        return True
    return False


def looks_like_data_table(text: str) -> bool:
    """Numeric/grid table rather than multi-panel marketing diagram or node list."""
    st = _ocr_line_stats(text)
    if st["n_lines"] < 3 or st["chars"] < 40:
        return False
    alnum = sum(1 for c in text if c.isalnum()) or 1
    digits = sum(1 for c in text if c.isdigit())
    spacey = text.count("  ") + text.count("\t")
    cidrs = len(re.findall(r"\d+\.\d+\.\d+\.\d+/\d+", text))
    # Explicit multi-row subnet/CIDR grids
    if cidrs >= 2 and st["n_lines"] >= 3:
        return True
    # Digit-heavy with column-like spacing (not just "name IP" labels)
    if digits > alnum * 0.22 and spacey >= 5 and st["n_lines"] >= 4:
        return True
    # Header-ish first line with multiple short numeric rows
    if st["n_lines"] >= 4 and spacey >= 3 and digits > alnum * 0.2:
        return True
    return False


def _normalize_for_match(s: str) -> str:
    """Lowercase, collapse non-alnum runs — for coverage checks."""
    s = (s or "").lower()
    s = re.sub(r"[^a-z0-9]+", " ", s)
    return re.sub(r"\s+", " ", s).strip()


def is_ocr_garble(token: str) -> bool:
    """True for broken OCR fragments we should never merge into vision prose."""
    t = (token or "").strip()
    if len(t) < 3:
        return True
    alnum = sum(1 for c in t if c.isalnum())
    if alnum < 3:
        return True
    # Too much punctuation / junk symbols
    junk = sum(1 for c in t if c in "|[]{}<>~`@#$%^*_+=")
    if junk >= 2:
        return True
    letters = sum(1 for c in t if c.isalpha())
    if letters >= 4:
        # Vowel-less noise (e.g. "NEE ee", "py Sa EN") — loose check
        vowels = sum(1 for c in t.lower() if c in "aeiou")
        if vowels == 0 and not re.search(r"\d", t):
            return True
    # Repeated word doubled by OCR ("Overlay Network Overlay Network")
    parts = t.split()
    if len(parts) >= 4 and parts[: len(parts) // 2] == parts[len(parts) // 2 :]:
        return True
    # Mid-word broken fragments with stray single letters between words
    if re.search(r"\b[a-zA-Z]\b.*\b[a-zA-Z]\b", t) and len(parts) >= 3:
        singles = sum(1 for p in parts if len(p) == 1)
        if singles >= 2:
            return True
    # Ends with OCR truncation markers
    if re.search(r"(in:|in;$|tion:$)", t, re.I) and len(t) < 28:
        # e.g. "Just-in-time provisionin:" — treat as garble if short truncated form
        if t.rstrip(":").lower().endswith(("in", "tion", "ment")):
            pass  # may still be useful if not covered — don't auto-garble all
    return False


# Short ALLCAPS only kept if known tech (avoids OCR junk like "SEU", "NEE")
_KNOWN_ACRONYMS = frozenset({
    "RBAC", "TACACS", "VRF", "ALB", "WAF", "VPN", "ACL", "API", "SSH", "TLS",
    "SSL", "DNS", "NTP", "AAA", "SSO", "MFA", "PCI", "AWS", "GCP", "SCP", "FTP",
    "S3", "SSP", "OMP", "TLOC", "IPSEC", "DTLS", "SAML", "HTTPS", "HTTP", "TCP",
    "UDP", "NAT", "BGP", "OSPF", "MPLS", "DMZ", "FQDN", "CIDR", "VLAN", "VXLAN",
    "GRE", "IKE", "PKI", "LDAP", "RADIUS", "SNMP", "MTU", "QOS", "SLA", "HA",
    "DR", "CI", "CD", "QA", "UI", "CLI", "SDK", "AMI", "VPC", "IAM",
})


def is_high_value_ocr_token(token: str) -> bool:
    """Only merge identifiers that matter for retrieval if vision missed them."""
    t = (token or "").strip()
    if not t or is_ocr_garble(t):
        return False
    # Pure URL / IP / CIDR tokens (whole string)
    if re.fullmatch(r"https?://\S+", t, re.I):
        return True
    if re.fullmatch(r"(?:\d{1,3}\.){3}\d{1,3}(?:/\d{1,2})?", t):
        return True
    if _URL_RE.search(t) and " " not in t:
        return True
    if _IP_RE.fullmatch(t) or re.fullmatch(r"\d+\.\d+\.\d+\.\d+/\d+", t):
        return True
    if _HOST_RE.search(t) and " " not in t:
        return True
    # Protocol or port-ish
    if re.fullmatch(r"(?:HTTPS?|TCP|UDP|SAML|TLS|DTLS|IPsec)(?:\s*\(?\d{2,5}\)?)?", t, re.I):
        return True
    # Acronyms: known set, or longer ALLCAPS (TACACS-length)
    if re.fullmatch(r"[A-Z]{2,12}", t):
        return t in _KNOWN_ACRONYMS or len(t) >= 5
    # Known product / platform tokens
    if re.search(
        r"\b(?:vManage|vBond|vSmart|vMonitor|vOrchestrator|vdaemon|SD-WAN|"
        r"Self-Service|Portal|Jenkins|Bamboo|SonarQube|Okta|TACACS|RBAC)\b",
        t,
        re.I,
    ):
        # Reject if mixed with obvious garble words
        if re.search(r"\b(?:py|ee|ted|in:|NEE|SEU|Pisco)\b", t, re.I):
            return False
        return True
    # Multi-word Title Case / product labels — every word capitalized, no junk
    words = t.replace(",", " ").split()
    if 2 <= len(words) <= 6:
        if all(re.match(r"^[A-Z0-9][A-Za-z0-9+./-]{0,30}$", w) for w in words):
            if sum(len(w) for w in words) >= 8 and not any(len(w) == 1 for w in words):
                return True
    # Single long technical token
    if re.fullmatch(r"[A-Za-z][A-Za-z0-9+./-]{5,40}", t) and t[0].isupper():
        return True
    return False


def ocr_token_covered_by_vision(token: str, vision_text: str, *, ratio: float = 0.82) -> bool:
    """True if vision already contains this token (exact, substring, or fuzzy)."""
    from difflib import SequenceMatcher

    if not token or not vision_text:
        return False
    vis_raw = vision_text.lower()
    tok_l = token.lower().strip()
    if tok_l in vis_raw:
        return True
    # Strip punctuation variants
    if re.sub(r"[^a-z0-9]+", "", tok_l) in re.sub(r"[^a-z0-9]+", "", vis_raw):
        return True
    tok_n = _normalize_for_match(token)
    vis_n = _normalize_for_match(vision_text)
    if not tok_n:
        return True
    if tok_n in vis_n:
        return True
    # Token words all present as whole words in vision
    words = [w for w in tok_n.split() if len(w) > 2]
    if words and all(re.search(rf"\b{re.escape(w)}\b", vis_n) for w in words):
        return True
    # Fuzzy: compare against vision lines / bullet snippets
    candidates = re.split(r"[\n,;|/]+", vision_text)
    candidates.append(vision_text)
    for cand in candidates:
        c_n = _normalize_for_match(cand)
        if not c_n or len(c_n) < 3:
            continue
        # Sliding window for long vision lines
        if len(c_n) <= max(80, len(tok_n) + 20):
            if SequenceMatcher(None, tok_n, c_n).ratio() >= ratio:
                return True
        elif tok_n in c_n:
            return True
        else:
            # local window
            for i in range(0, max(1, len(c_n) - len(tok_n) + 1), max(1, len(tok_n) // 2)):
                window = c_n[i : i + len(tok_n) + 8]
                if SequenceMatcher(None, tok_n, window).ratio() >= ratio:
                    return True
    return False


def extract_ocr_tokens(text: str) -> list[str]:
    """Candidate OCR strings (pre-filter); merge applies high-value + coverage gates.

    Extract high-precision patterns first (URLs, IPs) so they are not swallowed
    by broader Title-Case matches.
    """
    if not text:
        return []
    found: list[str] = []
    seen: set[str] = set()

    def _add(tok: str) -> None:
        tok = tok.strip().rstrip(".,;:)]}")
        key = tok.lower()
        if len(tok) < 3 or key in seen:
            return
        if sum(1 for c in tok if c.isalnum()) < 2:
            return
        if is_ocr_garble(tok):
            return
        seen.add(key)
        found.append(tok)

    # Priority 1: full URLs and IPs/CIDRs
    for m in _URL_RE.finditer(text):
        _add(m.group(0))
    for m in re.finditer(r"\b(?:\d{1,3}\.){3}\d{1,3}(?:/\d{1,2})?\b", text):
        _add(m.group(0))
    # Priority 2: other regex tokens (skip if already seen)
    for m in _OCR_TOKEN_RE.finditer(text):
        _add(m.group(0))
    # Priority 3: clean label lines (no lowercase junk words)
    for ln in _ocr_line_stats(text)["lines"]:
        if 3 <= len(ln) <= 48 and ln[0].isupper() and not ln.endswith("."):
            if re.search(r"\b(?:py|ee|ted|NEE|Pisco|selese)\b", ln, re.I):
                continue
            if re.search(r"[A-Za-z]{3,}", ln):
                _add(ln)
    return found[:80]


def select_ocr_tokens_to_merge(vision_text: str, ocr_text: str, *, max_items: int = 10) -> list[str]:
    """Phase A: high-value OCR identifiers not already covered by vision."""
    body = vision_text or ""
    keep: list[str] = []
    for tok in extract_ocr_tokens(ocr_text or ""):
        if not is_high_value_ocr_token(tok):
            continue
        if ocr_token_covered_by_vision(tok, body):
            continue
        keep.append(tok)
        if len(keep) >= max_items:
            break
    return keep


def merge_ocr_into_vision(vision_text: str, ocr_text: str, page: int, kind: str) -> str:
    """Append only high-value OCR identifiers missing from vision (Phase A filter).

    When vision already did a good job, this returns vision prose with no OCR
    appendix — no garble dump.
    """
    body = _wrap_generic(page, kind, vision_text or "")
    missing = select_ocr_tokens_to_merge(body, ocr_text or "")
    if not missing:
        return body
    extra = (
        "Also visible (OCR identifiers not in description):\n"
        + "\n".join(f"- {t}" for t in missing)
    )
    if "[/Figure]" in body:
        return body.replace("[/Figure]", extra + "\n[/Figure]", 1)
    return body + "\n" + extra


def classify_region(region: ImageRegion) -> ImageKind:
    """Heuristic classifier: prefer vision for diagrams; OCR-only for linear prose."""
    text = region.ocr_text or ""
    chars = len(text)
    st = _ocr_line_stats(text)
    alnum = sum(1 for c in text if c.isalnum()) or 1
    digits = sum(1 for c in text if c.isdigit())
    ip_hits = len(_IP_RE.findall(text))
    low = text.lower()

    # 1) Data tables (numeric grids) — not multi-panel architecture art
    if looks_like_data_table(text):
        return ImageKind.TABLE

    # 2) Linear prose / terminal dump — OCR-only is appropriate
    if looks_like_linear_prose(text) and not looks_like_diagram(region, text):
        return ImageKind.TEXT_SCAN

    # 3) Spatial diagrams — NEVER text_scan even if OCR is wordy (p.22 case)
    if looks_like_diagram(region, text):
        process_hits = sum(1 for k in _PROCESS_HINTS if k in low)
        network_hits = sum(1 for k in _NETWORK_HINTS if k in low)
        if ip_hits >= 2 or _HOST_RE.search(text) or network_hits > process_hits:
            return ImageKind.NETWORK
        return ImageKind.PROCESS

    # 4) Host/IP label soup on medium+ images
    if (ip_hits >= 2 or _HOST_RE.search(text)) and region.width >= 400:
        return ImageKind.NETWORK

    # 5) Sparse OCR — large images without diagram keywords are often
    # illustrations/photos (still describe, via OTHER/illustration path)
    if chars < 25:
        if any(w in low for w in ("start", "end", "decision", "workflow", "approve")):
            return ImageKind.PROCESS
        if region.source == "page_render" and looks_like_diagram(region, text):
            return ImageKind.NETWORK
        return ImageKind.OTHER

    # 6) Remaining mid-length OCR on large canvas with short lines → diagram
    if region.width >= 600 and region.height >= 250 and st["n_lines"] >= 3 and st["avg_len"] < 55:
        return ImageKind.PROCESS if digits < alnum * 0.12 else ImageKind.NETWORK

    return ImageKind.OTHER


# ---------------------------------------------------------------------------
# Markdown / CSV table parsing (strategy TABLE)
# ---------------------------------------------------------------------------

def parse_markdown_table(text: str) -> list[list[str]] | None:
    """Parse a GFM markdown table or simple CSV into rows. None if not a table."""
    if not text:
        return None
    cleaned = text.strip()
    if re.search(r"\bNOT_A_TABLE\b", cleaned, re.I):
        return None
    # Strip fences
    cleaned = re.sub(r"```(?:markdown|md|csv|text)?\s*", "", cleaned)
    cleaned = cleaned.replace("```", "").strip()

    lines = [ln.strip() for ln in cleaned.splitlines() if ln.strip()]
    # Markdown table lines
    md_lines = [ln for ln in lines if "|" in ln]
    if len(md_lines) >= 2:
        rows: list[list[str]] = []
        for ln in md_lines:
            # Skip separator |---|---|
            if re.match(r"^\|?[\s:-]+\|[\s|:-]+$", ln):
                continue
            cells = [c.strip() for c in ln.strip("|").split("|")]
            if cells:
                rows.append(cells)
        if len(rows) >= 2:
            return rows

    # CSV fallback (comma-separated, ≥2 lines, ≥2 cols)
    csv_rows: list[list[str]] = []
    for ln in lines:
        if "," not in ln:
            continue
        # naive CSV split
        cells = [c.strip().strip('"') for c in ln.split(",")]
        if len(cells) >= 2:
            csv_rows.append(cells)
    if len(csv_rows) >= 2:
        return csv_rows
    return None


# ---------------------------------------------------------------------------
# Vision prompts
# ---------------------------------------------------------------------------

_SYSTEM = (
    "You are a precise document figure extractor for a search index. "
    "Always produce a useful description of what is in the image so readers can find it later. "
    "Never invent IP addresses, hostnames, ports, labels, or cell values. "
    "If text is unreadable, say so rather than guessing. "
    "Do not refuse: if the image is non-technical (photo, illustration, logo, artwork), "
    "still describe it clearly and completely."
)

_PROMPT_TABLE = """This image is a table from a technical document.
Output a GitHub-flavored markdown table only.
Rules:
- Transcribe every visible cell exactly (IPs, hostnames, ports, codes, numbers).
- Keep row/column alignment; one markdown row per table row.
- Do not invent, round, or fill missing cells — use empty cells.
- If it is not a table, reply exactly: NOT_A_TABLE
"""

_PROMPT_NETWORK = """This image may be a network / infrastructure diagram from a document.
If it IS a network/topology/wiring diagram, use exactly this structure:

[Figure p.{page} — network]
Identifiers:
- (list every visible IP, hostname, FQDN, ASN, VRF, interface, port — exact text only)
Nodes:
- (components / devices / services)
Links:
- A --protocol/port--> B  (only if shown)
Zones:
- (trust boundaries, clouds, sites if shown)
Notes:
- (other visible labels; do not invent)
[/Figure]

If it is NOT a technical network diagram (photo, illustration, logo, artwork, etc.),
do NOT refuse. Instead describe it fully for search indexing:

[Figure p.{page} — illustration]
Subject:
- (what/who is depicted)
Setting / scene:
- ...
Objects / characters:
- ...
Visible text:
- (any captions, watermarks, labels — exact)
Style:
- (photo, cartoon, watercolor, logo, icon, etc.)
Notes:
- (colors, mood, notable details)
[/Figure]
"""

_PROMPT_PROCESS = """This image may be a process, architecture, portal, or multi-panel diagram
(flowchart, lifecycle, benefits graphic, org chart, CI/CD, etc.).
If it IS that kind of technical/business diagram, use exactly this structure:

[Figure p.{page} — process]
Layout:
- (center title / main hub if any; left/right/top/bottom panels)
Components / pillars:
- Name: bullet points exactly as shown
- ...
Flows / relationships:
- A --> B (only if shown)
Identifiers:
- (every product name, URL, port, protocol, acronym — exact text)
Notes:
- ...
[/Figure]

If it is NOT a technical diagram (photo, illustration, logo, clip art, artwork, etc.),
do NOT refuse and do NOT say you cannot fulfill the request. Describe it fully:

[Figure p.{page} — illustration]
Subject:
- (what/who is depicted)
Setting / scene:
- ...
Objects / characters:
- ...
Visible text:
- (any captions, watermarks, brand names — exact)
Style:
- (photo, cartoon, watercolor, logo, icon, etc.)
Notes:
- (colors, mood, notable details useful for search)
[/Figure]

Rules:
- Transcribe ALL visible labels exactly (do not invent).
- If an OCR hint is provided, include every OCR string you can confirm is visible.
"""

_PROMPT_ILLUSTRATION = """Describe this document image in detail for a full-text search index.
It may be a photo, illustration, logo, icon, screenshot, or diagram.
Always provide a rich description — never refuse.

[Figure p.{page} — illustration]
Subject:
- (main subject)
Setting / scene:
- ...
Objects / characters:
- ...
Visible text:
- (every readable word/label exactly; none if blank)
Style:
- (photo, cartoon, watercolor, logo, UI screenshot, etc.)
Colors / composition:
- ...
Notes:
- (anything someone might search for)
[/Figure]
"""

_PROMPT_OTHER = """Describe this document image in detail for search indexing.
- Transcribe ALL visible text exactly (labels, legends, IPs, hostnames, URLs, brand names).
- Describe what is shown even if non-technical (photos, illustrations, logos).
- Do not invent text or identifiers that are not visible.
- Never refuse to describe the image.

Wrap the answer as:
[Figure p.{page} — figure]
...detailed description...
[/Figure]
"""

_PROMPT_CLASSIFY = """Classify this document image. Reply with exactly one word:
table
network
process
text_scan
other

Use text_scan ONLY for linear paragraph text or terminal dumps.
Use process for multi-panel architecture, portal, lifecycle, or flowchart graphics.
Use network for topology / wiring diagrams.
"""


def _document_context_hint(region: ImageRegion) -> str:
    """Give vision local context without claiming that it is visible pixels."""
    fields: list[str] = []
    if region.section_path:
        fields.append("Section: " + " > ".join(region.section_path))
    if region.caption:
        fields.append(f"Caption: {region.caption}")
    if region.alt_text:
        fields.append(f"Office alt text: {region.alt_text}")
    if region.previous_text:
        fields.append(f"Text immediately before: {region.previous_text[:1000]}")
    if region.following_text:
        fields.append(f"Text immediately after: {region.following_text[:1000]}")
    if not fields:
        return ""
    return (
        "\n\nDOCUMENT CONTEXT (use to disambiguate the image, but do not report "
        "context text as visually observed unless you can confirm it in the image):\n"
        + "\n".join(fields)
    )


def _format_figure_for_placement(region: ImageRegion, text: str) -> str:
    """Retag model prose with a stable figure ID and source context."""
    clean = (text or "").strip()
    if not clean:
        return ""
    label = region.figure_id or f"p.{region.page + 1}"
    clean = re.sub(r"^\[Figure\s+[^\]—]+\s*—", f"[Figure {label} —", clean, count=1)
    if not clean.startswith("[Figure"):
        clean = f"[Figure {label} — {region.kind.value}]\n{clean}\n[/Figure]"

    provenance: list[str] = []
    if region.section_path:
        provenance.append("Section: " + " > ".join(region.section_path))
    if region.caption:
        provenance.append(f"Caption: {region.caption}")
    if region.alt_text:
        provenance.append(f"Alt text: {region.alt_text}")
    if provenance:
        first, sep, rest = clean.partition("\n")
        clean = first + "\n" + "\n".join(provenance) + (sep + rest if sep else "")
    return clean


def _vision(image_bytes: bytes, user_prompt: str, max_tokens: int | None = None) -> str:
    from src.generation.llm_client import generate_vision

    result = generate_vision(
        system_prompt=_SYSTEM,
        user_prompt=user_prompt,
        image_bytes=image_bytes,
        mime_type="image/png",
        temperature=0.0,
        max_tokens=max_tokens or settings.figure_vision_max_tokens,
        timeout=settings.figure_vision_timeout_seconds,
    )
    # Log a preview so operators can verify figure extraction quality live
    preview = (result or "").replace("\n", " \\n ")[:400]
    logger.info(f"Vision response ({len(result or '')} chars): {preview}")
    return result


def _wrap_ocr_block(page: int, text: str) -> str:
    return (
        f"[Figure p.{page + 1} — ocr]\n"
        f"{text.strip()}\n"
        f"[/Figure]"
    )


def _wrap_generic(page: int, kind: str, body: str) -> str:
    body = body.strip()
    if body.startswith("[Figure"):
        return body
    return f"[Figure p.{page + 1} — {kind}]\n{body}\n[/Figure]"


# ---------------------------------------------------------------------------
# Strategies
# ---------------------------------------------------------------------------

def strategy_ocr(region: ImageRegion) -> tuple[list[SheetGrid], list[ProseBlock]]:
    text = region.ocr_text or ocr_image(region.image_bytes)
    if len(text.strip()) < 15:
        return [], []
    block = ProseBlock(text=_wrap_ocr_block(region.page, text), page=region.page)
    return [], [block]


def strategy_table(region: ImageRegion) -> tuple[list[SheetGrid], list[ProseBlock]]:
    hint = ""
    if region.ocr_text and len(region.ocr_text) > 20:
        hint = f"\n\nOCR hint (may be noisy; prefer what you see in the image):\n{region.ocr_text[:1500]}"
    try:
        raw = _vision(
            region.image_bytes,
            _PROMPT_TABLE + hint + _document_context_hint(region),
            max_tokens=3000,
        )
    except Exception as e:
        logger.warning(f"TABLE vision failed p{region.page} img{region.index}: {e}")
        return strategy_ocr(region)

    rows = parse_markdown_table(raw)
    if not rows:
        logger.info(f"TABLE vision not a table p{region.page} img{region.index}; falling back")
        # Architecture graphics often mis-route to TABLE; use process vision + OCR merge
        if looks_like_diagram(region, region.ocr_text or ""):
            return strategy_process(region)
        if raw and "NOT_A_TABLE" not in raw.upper() and len(raw) > 40:
            merged = merge_ocr_into_vision(raw, region.ocr_text, region.page, "figure")
            return [], [ProseBlock(text=merged, page=region.page)]
        return strategy_ocr(region)

    name = f"p{region.page}_imgtable{region.index}"
    grid = normalize_grid(rows, sheet_name=name)
    if not grid.rows or len(grid.rows) < 2:
        return strategy_ocr(region)

    # Short topic blurb for vector search (not the full grid)
    blurb = (
        f"[Figure p.{region.page + 1} — table]\n"
        f"Structured table extracted as '{name}' "
        f"({len(grid.rows)} rows × {len(grid.rows[0]) if grid.rows else 0} cols). "
        f"Queryable via structured store.\n"
        f"[/Figure]"
    )
    return [grid], [ProseBlock(text=blurb, page=region.page)]


def _ocr_hint_block(region: ImageRegion, limit: int = 1500) -> str:
    if not region.ocr_text or len(region.ocr_text) < 15:
        return ""
    tokens = extract_ocr_tokens(region.ocr_text)
    parts = [f"\n\nOCR hint (noisy; include any you can confirm visible):\n{region.ocr_text[:limit]}"]
    if tokens:
        parts.append("OCR tokens checklist:\n" + "\n".join(f"- {t}" for t in tokens[:40]))
    return "".join(parts)


_REFUSAL_RE = re.compile(
    r"(?i)\b("
    r"cannot fulfill|can'?t fulfill|unable to (?:provide|describe|fulfill)|"
    r"does not contain (?:a |any )?(?:technical|process|architecture)|"
    r"not a technical|not a process|not a network|does not meet the criteria|"
    r"i cannot|i'?m unable|i am unable"
    r")\b"
)


def looks_like_refusal(text: str) -> bool:
    """True when the model declined to describe instead of indexing the image."""
    t = (text or "").strip()
    if not t:
        return True
    if _REFUSAL_RE.search(t):
        # Real descriptions that mention "not a technical diagram" then still
        # describe: require refusal to dominate short answers
        if len(t) < 500:
            return True
        # Long answer that starts with refusal but continues with illustration —
        # keep if it already has a usable illustration block
        if "[Figure" in t and "illustration" in t.lower() and len(t) > 400:
            return False
        if t.lower().count("cannot") + t.lower().count("unable") >= 1 and len(t) < 700:
            return True
    return False


def strategy_illustration(region: ImageRegion) -> tuple[list[SheetGrid], list[ProseBlock]]:
    """Always-on descriptive path for photos, art, logos, and soft content."""
    prompt = (
        _PROMPT_ILLUSTRATION.format(page=region.page + 1)
        + _ocr_hint_block(region)
        + _document_context_hint(region)
    )
    try:
        raw = _vision(region.image_bytes, prompt)
    except Exception as e:
        logger.warning(f"ILLUSTRATION vision failed p{region.page} img{region.index}: {e}")
        return strategy_ocr(region)
    if len(raw.strip()) < 20:
        return strategy_ocr(region)
    # Strip leading refusal sentences if any slipped through
    cleaned = raw.strip()
    if looks_like_refusal(cleaned) and "[Figure" not in cleaned:
        # Last resort: wrap a minimal note so we still keep something
        cleaned = (
            f"[Figure p.{region.page + 1} — illustration]\n"
            f"Subject:\n- (model returned a non-descriptive reply)\n"
            f"Notes:\n- {cleaned[:300]}\n"
            f"[/Figure]"
        )
    merged = merge_ocr_into_vision(cleaned, region.ocr_text, region.page, "illustration")
    return [], [ProseBlock(text=merged, page=region.page)]


def _vision_with_illustration_fallback(
    region: ImageRegion,
    primary_prompt: str,
    *,
    tag: str,
) -> tuple[list[SheetGrid], list[ProseBlock]]:
    try:
        raw = _vision(region.image_bytes, primary_prompt)
    except Exception as e:
        logger.warning(f"{tag} vision failed p{region.page} img{region.index}: {e}")
        return strategy_illustration(region)
    if len(raw.strip()) < 20 or looks_like_refusal(raw):
        logger.info(
            f"figure p{region.page} img{region.index}: {tag} refused/empty → illustration describe"
        )
        return strategy_illustration(region)
    # Prefer illustration tag if model already switched formats
    kind = "illustration" if "— illustration]" in raw or "— logo]" in raw else tag.lower()
    merged = merge_ocr_into_vision(raw, region.ocr_text, region.page, kind)
    return [], [ProseBlock(text=merged, page=region.page)]


def strategy_network(region: ImageRegion) -> tuple[list[SheetGrid], list[ProseBlock]]:
    prompt = (
        _PROMPT_NETWORK.format(page=region.page + 1)
        + _ocr_hint_block(region)
        + _document_context_hint(region)
    )
    return _vision_with_illustration_fallback(region, prompt, tag="network")


def strategy_process(region: ImageRegion) -> tuple[list[SheetGrid], list[ProseBlock]]:
    prompt = (
        _PROMPT_PROCESS.format(page=region.page + 1)
        + _ocr_hint_block(region)
        + _document_context_hint(region)
    )
    return _vision_with_illustration_fallback(region, prompt, tag="process")


def strategy_other(region: ImageRegion) -> tuple[list[SheetGrid], list[ProseBlock]]:
    # Prefer dedicated illustration prompt for richer non-tech descriptions
    return strategy_illustration(region)


def _vision_classify(region: ImageRegion) -> ImageKind | None:
    try:
        label = _vision(region.image_bytes, _PROMPT_CLASSIFY, max_tokens=16).strip().lower()
        label = label.split()[0] if label else ""
        label = label.strip(".,:;\"'")
        mapping = {
            "table": ImageKind.TABLE,
            "network": ImageKind.NETWORK,
            "process": ImageKind.PROCESS,
            "text_scan": ImageKind.TEXT_SCAN,
            "text": ImageKind.TEXT_SCAN,
            "ocr": ImageKind.TEXT_SCAN,
            "other": ImageKind.OTHER,
            "diagram": ImageKind.PROCESS,
            "flowchart": ImageKind.PROCESS,
            "architecture": ImageKind.PROCESS,
            "screenshot": ImageKind.OTHER,
        }
        return mapping.get(label)
    except Exception as e:
        logger.debug(f"vision classify failed: {e}")
        return None


def run_strategy(region: ImageRegion) -> tuple[list[SheetGrid], list[ProseBlock]]:
    kind = region.kind
    if kind == ImageKind.SKIP:
        return [], []
    if kind == ImageKind.TEXT_SCAN:
        # Safety net: never OCR-only a large diagram that was mis-labeled
        if looks_like_diagram(region, region.ocr_text or ""):
            logger.info(
                f"figure p{region.page} img{region.index}: upgrading text_scan → process (diagram-like)"
            )
            region.kind = ImageKind.PROCESS
            return strategy_process(region)
        return strategy_ocr(region)
    if kind == ImageKind.TABLE:
        return strategy_table(region)
    if kind == ImageKind.NETWORK:
        return strategy_network(region)
    if kind == ImageKind.PROCESS:
        return strategy_process(region)
    # OTHER: technical-looking layout → process; else full illustration describe
    if looks_like_diagram(region, region.ocr_text or ""):
        return strategy_process(region)
    if region.ocr_text and len(region.ocr_text) > 200 and looks_like_linear_prose(region.ocr_text):
        return strategy_ocr(region)
    return strategy_illustration(region)


# ---------------------------------------------------------------------------
# Merge helpers
# ---------------------------------------------------------------------------

def merge_prose_by_page(
    base: list[ProseBlock],
    extra: list[ProseBlock],
) -> list[ProseBlock]:
    """Interleave prose and figure blocks by page number (stable within page)."""
    combined = list(base) + list(extra)
    # sort by page, then keep relative order via original index
    indexed = list(enumerate(combined))
    indexed.sort(key=lambda iv: (iv[1].page, iv[0]))
    return [b for _, b in indexed]


def _looks_like_pdf_heading(block: ProseBlock, page_blocks: list[ProseBlock]) -> bool:
    text = block.text.strip()
    if not text or len(text) > 180:
        return False
    sizes = sorted(b.font_size for b in page_blocks if b.font_size and b.font_size > 0)
    median = sizes[len(sizes) // 2] if sizes else 0
    if block.font_size and median and block.font_size >= median * 1.15:
        return True
    return bool(re.match(r"(?i)^(section|chapter|appendix)\s+[\w.:-]+", text))


def attach_pdf_context(
    regions: list[ImageRegion],
    extracted: ExtractedPdf,
) -> None:
    """Attach nearby positioned PDF text before vision analysis."""
    layout_by_page: dict[int, list[ProseBlock]] = {}
    prose_by_page: dict[int, list[ProseBlock]] = {}
    for block in extracted.layout_blocks:
        layout_by_page.setdefault(block.page, []).append(block)
    for block in extracted.prose_blocks:
        prose_by_page.setdefault(block.page, []).append(block)

    for region in regions:
        blocks = sorted(
            layout_by_page.get(region.page, []),
            key=lambda b: (b.bbox[1], b.bbox[0]) if b.bbox else (0, 0),
        )
        if not blocks or not region.bbox or region.source == "page_render":
            page_text = "\n".join(b.text for b in prose_by_page.get(region.page, []))
            region.following_text = page_text[:1200]
            headings = [b.text for b in blocks if _looks_like_pdf_heading(b, blocks)]
            if headings:
                region.section_path = [headings[0]]
            continue

        _, top, _, bottom = region.bbox
        before = [
            b for b in blocks if b.bbox and b.bbox[3] <= top + 3
        ]
        after = [
            b for b in blocks if b.bbox and b.bbox[1] >= bottom - 3
        ]
        region.previous_text = "\n".join(b.text for b in before[-3:])[-1200:]
        region.following_text = "\n".join(b.text for b in after[:3])[:1200]

        caption_candidates = before[-2:] + after[:2]
        caption_candidates.sort(key=lambda b: min(
            abs((b.bbox[3] if b.bbox else top) - top),
            abs((b.bbox[1] if b.bbox else bottom) - bottom),
        ))
        for candidate in caption_candidates:
            if re.match(r"(?i)^\s*(figure|fig\.)\s*\d+", candidate.text):
                region.caption = candidate.text.strip()
                break

        headings = [
            b for b in before if _looks_like_pdf_heading(b, blocks)
        ]
        if headings:
            region.section_path = [headings[-1].text.strip()]


def merge_pdf_prose_by_position(
    extracted: ExtractedPdf,
    figures: list[FigureRecord],
) -> list[ProseBlock]:
    """Interleave digital PDF text and figures using top-left page coordinates."""
    base_by_page: dict[int, list[ProseBlock]] = {}
    layout_by_page: dict[int, list[ProseBlock]] = {}
    figures_by_page: dict[int, list[FigureRecord]] = {}
    for block in extracted.prose_blocks:
        base_by_page.setdefault(block.page, []).append(block)
    for block in extracted.layout_blocks:
        layout_by_page.setdefault(block.page, []).append(block)
    for record in figures:
        if record.page is not None:
            figures_by_page.setdefault(record.page, []).append(record)

    pages = sorted(set(base_by_page) | set(layout_by_page) | set(figures_by_page))
    merged: list[ProseBlock] = []
    for page in pages:
        page_figures = figures_by_page.get(page, [])
        page_layout = layout_by_page.get(page, [])
        if not page_figures:
            merged.extend(base_by_page.get(page, []))
            continue
        if not page_layout:
            # Scanned/full-page figures have page-level, not object-level,
            # placement. Put the page description before OCR prose.
            for record in page_figures:
                merged.append(ProseBlock(
                    text=record.description, page=page, bbox=record.bbox,
                    content_type="figure", figure_id=record.figure_id,
                ))
            merged.extend(base_by_page.get(page, []))
            continue

        events: list[tuple[float, float, str, object]] = []
        for block in page_layout:
            if not block.bbox:
                continue
            cx = (block.bbox[0] + block.bbox[2]) / 2
            cy = (block.bbox[1] + block.bbox[3]) / 2
            inside_figure = any(
                record.bbox
                and record.source == "embedded"
                and record.bbox[0] <= cx <= record.bbox[2]
                and record.bbox[1] <= cy <= record.bbox[3]
                for record in page_figures
            )
            if not inside_figure:
                events.append((block.bbox[1], block.bbox[0], "text", block))
        for record in page_figures:
            top = record.bbox[1] if record.bbox else float("inf")
            left = record.bbox[0] if record.bbox else float("inf")
            events.append((top, left, "figure", record))
        events.sort(key=lambda event: (event[0], event[1], event[2] != "figure"))

        pending_text: list[str] = []

        def flush_text() -> None:
            if pending_text:
                merged.append(ProseBlock(text="\n".join(pending_text), page=page))
                pending_text.clear()

        for _, _, event_type, value in events:
            if event_type == "text":
                pending_text.append(value.text)
            else:
                flush_text()
                record = value
                merged.append(ProseBlock(
                    text=record.description, page=page, bbox=record.bbox,
                    content_type="figure", figure_id=record.figure_id,
                ))
        flush_text()
    return merged


def pages_with_digital_tables(extracted: ExtractedPdf) -> set[int]:
    """Page indices that already have pdfplumber tables (avoid double table extract)."""
    pages: set[int] = set()
    for g in extracted.table_grids or []:
        # sheet names like p0_table0
        m = re.match(r"p(\d+)_", g.sheet_name or "")
        if m:
            pages.add(int(m.group(1)))
    return pages


# ---------------------------------------------------------------------------
# Public entry points
# ---------------------------------------------------------------------------

def process_image_regions(
    regions: list[ImageRegion],
    *,
    path_name: str = "",
    digital_table_pages: set[int] | None = None,
    progress_cb=None,
) -> FigureEnrichmentResult:
    """Run classify + strategies on regions. Shared by PDF and Office paths."""
    result = FigureEnrichmentResult(figures_seen=len(regions))
    if not regions:
        return result

    digital_table_pages = digital_table_pages or set()
    max_n = max(0, settings.figure_max_per_doc)
    regions = sorted(regions, key=lambda r: (r.page, r.index))
    # The budget applies to distinct image content, not placements. Repeated
    # figures reuse analysis but remain present everywhere they occur.
    selected: list[ImageRegion] = []
    selected_hashes: set[str] = set()
    for region in regions:
        digest = region.content_hash or f"placement:{region.page}:{region.index}"
        if digest in selected_hashes:
            selected.append(region)
            continue
        if max_n and len(selected_hashes) >= max_n:
            result.figures_skipped += 1
            continue
        selected_hashes.add(digest)
        selected.append(region)
    regions = selected
    analysis_cache: dict[
        str, tuple[ImageKind, str, list[SheetGrid], list[ProseBlock]]
    ] = {}

    for i, region in enumerate(regions):
        if progress_cb:
            try:
                progress_cb(f"Extracting figures ({i + 1}/{len(regions)})…")
            except Exception:
                pass

        cached = analysis_cache.get(region.content_hash) if region.content_hash else None
        if cached:
            kind, cached_ocr, cached_grids, cached_prose = cached
            region.ocr_text = cached_ocr
            # Structured image tables are stored once; prose is placed at every
            # occurrence with placement-specific context below.
            grids, prose = [], list(cached_prose)
        else:
            if settings.figure_ocr_first:
                region.ocr_text = ocr_image(region.image_bytes)

            kind = classify_region(region)
            if kind == ImageKind.OTHER and region.width * region.height > 80_000:
                vkind = _vision_classify(region)
                if vkind is not None and vkind != ImageKind.TEXT_SCAN:
                    kind = vkind
                elif vkind == ImageKind.TEXT_SCAN and looks_like_diagram(region, region.ocr_text or ""):
                    kind = ImageKind.PROCESS

        if (
            kind == ImageKind.TABLE
            and region.source == "page_render"
            and region.page in digital_table_pages
        ):
            logger.info(
                f"figure extract: skip page_render TABLE on p{region.page} "
                f"(digital tables already present)"
            )
            result.figures_skipped += 1
            continue

        region.kind = kind
        ocr_preview = (region.ocr_text or "").replace("\n", " ")[:120]
        logger.info(
            f"figure p{region.page} img{region.index}: kind={kind.value} "
            f"source={region.source} size={region.width}x{region.height} "
            f"bytes={len(region.image_bytes)} ocr_preview={ocr_preview!r}"
        )
        if not cached:
            try:
                grids, prose = run_strategy(region)
            except Exception as e:
                logger.warning(f"figure strategy failed p{region.page} img{region.index}: {e}")
                result.figures_skipped += 1
                continue
            if region.content_hash:
                analysis_cache[region.content_hash] = (
                    region.kind, region.ocr_text, list(grids), list(prose),
                )

        placed_prose = [
            ProseBlock(text=_format_figure_for_placement(region, b.text), page=b.page)
            for b in prose if b.text.strip()
        ]

        if grids or placed_prose:
            for g in grids:
                # Namespace sheet names by source for Office docs
                if region.source in ("docx", "xlsx", "pptx") and not g.sheet_name.startswith(region.source):
                    g.sheet_name = f"{region.source}_{g.sheet_name}"
                logger.info(
                    f"figure p{region.page} img{region.index}: TABLE grid "
                    f"{g.sheet_name} rows={len(g.rows)} cols={len(g.rows[0]) if g.rows else 0}"
                )
            for b in placed_prose:
                pprev = b.text.replace("\n", " \\n ")[:350]
                logger.info(
                    f"figure p{region.page} img{region.index}: prose ({len(b.text)} chars): {pprev}"
                )
            result.table_grids.extend(grids)
            result.prose_blocks.extend(placed_prose)
            description = "\n\n".join(b.text for b in placed_prose if b.text.strip())
            if description:
                result.figure_records.append(FigureRecord(
                    figure_id=region.figure_id or f"p{region.page + 1}-img{region.index + 1}",
                    description=description,
                    kind=region.kind.value,
                    content_hash=region.content_hash,
                    body_index=region.body_index,
                    section_path=list(region.section_path),
                    caption=region.caption,
                    alt_text=region.alt_text,
                    previous_text=region.previous_text,
                    following_text=region.following_text,
                    page=(
                        region.page
                        if region.source in ("embedded", "page_render")
                        else None
                    ),
                    bbox=region.bbox,
                    source=region.source,
                    slide=(region.page if region.source == "pptx" else None),
                ))
            result.figures_used += 1
        else:
            result.figures_skipped += 1

    logger.info(
        f"figure extract: {path_name or 'document'}: regions={result.figures_seen} "
        f"used={result.figures_used} skipped={result.figures_skipped} "
        f"grids={len(result.table_grids)} prose={len(result.prose_blocks)}"
    )
    return result


def enrich_pdf_with_figures(
    path: Path | str,
    extracted: ExtractedPdf,
    *,
    progress_cb=None,
) -> ExtractedPdf:
    """Run figure extraction and return a new ExtractedPdf with grids/prose merged.

    ``progress_cb(msg: str)`` optional status callback for the ingest queue UI.
    Fail-open: returns the original extracted content plus any successful figures.
    """
    if not settings.figure_extraction_enabled:
        return extracted

    path = Path(path)
    try:
        regions = extract_image_regions(path)
    except Exception as e:
        logger.warning(f"figure extract: region scan failed for {path.name}: {e}")
        return extracted

    if not regions:
        logger.info(f"figure extract: no image regions in {path.name}")
        return extracted

    attach_pdf_context(regions, extracted)
    enrich = process_image_regions(
        regions,
        path_name=path.name,
        digital_table_pages=pages_with_digital_tables(extracted),
        progress_cb=progress_cb,
    )
    if not enrich.table_grids and not enrich.prose_blocks:
        return extracted

    merged_prose = (
        merge_pdf_prose_by_position(extracted, enrich.figure_records)
        if enrich.figure_records
        else merge_prose_by_page(extracted.prose_blocks, enrich.prose_blocks)
    )
    existing_names = {g.sheet_name for g in extracted.table_grids}
    extra_grids = list(enrich.table_grids)
    for g in extra_grids:
        base = g.sheet_name
        n = 1
        while g.sheet_name in existing_names:
            g.sheet_name = f"{base}_{n}"
            n += 1
        existing_names.add(g.sheet_name)

    method = extracted.method
    if enrich.table_grids or enrich.prose_blocks:
        if method == "digital":
            method = "mixed"

    return ExtractedPdf(
        prose_blocks=merged_prose,
        table_grids=list(extracted.table_grids) + extra_grids,
        method=method,
        layout_blocks=list(extracted.layout_blocks),
        figure_records=list(enrich.figure_records),
    )


def enrich_office_document_with_figures(
    path: Path | str,
    base_text: str,
    *,
    document_blocks=None,
    progress_cb=None,
) -> OfficeFigureResult:
    """Return ordered Office prose, structured grids, and figure records.

    DOCX and PPTX descriptions are rendered at their source anchors. Other Office
    formats retain the figure appendix until they gain an ordered block parser.
    """
    if not settings.figure_extraction_enabled:
        return OfficeFigureResult(enriched_text=base_text or "")

    path = Path(path)
    try:
        regions = extract_image_regions_for_path(path)
    except Exception as e:
        logger.warning(f"figure extract: Office region scan failed for {path.name}: {e}")
        return OfficeFigureResult(enriched_text=base_text or "")

    if not regions:
        logger.info(f"figure extract: no image regions in {path.name}")
        return OfficeFigureResult(enriched_text=base_text or "")

    enrich = process_image_regions(
        regions, path_name=path.name, progress_cb=progress_cb,
    )
    figure_text = "\n\n".join(b.text for b in enrich.prose_blocks if b.text.strip())
    body = (base_text or "").strip()
    if path.suffix.lower() in (".docx", ".pptx") and enrich.figure_records:
        try:
            from src.ingestion.parser import parse_document, render_document_blocks

            blocks = document_blocks
            if blocks is None:
                blocks = parse_document(path).blocks
            by_id = {record.figure_id: record.description for record in enrich.figure_records}
            enriched_text = render_document_blocks(blocks, by_id)
        except Exception as e:
            logger.warning(f"figure extract: ordered Office render failed: {e}")
            enriched_text = body + ("\n\n## Embedded figures\n\n" + figure_text if body else figure_text)
    elif figure_text and body:
        enriched = body + "\n\n## Embedded figures\n\n" + figure_text
        enriched_text = enriched
    elif figure_text:
        enriched_text = figure_text
    else:
        enriched_text = body
    return OfficeFigureResult(
        enriched_text=enriched_text,
        table_grids=list(enrich.table_grids),
        figures=list(enrich.figure_records),
    )


def enrich_text_with_figures(
    path: Path | str,
    base_text: str,
    *,
    document_blocks=None,
    progress_cb=None,
) -> tuple[str, list[SheetGrid]]:
    """Compatibility wrapper returning the historical two-tuple."""
    result = enrich_office_document_with_figures(
        path,
        base_text,
        document_blocks=document_blocks,
        progress_cb=progress_cb,
    )
    return result.enriched_text, result.table_grids


async def enrich_pdf_with_figures_async(
    path: Path | str,
    extracted: ExtractedPdf,
    *,
    progress_cb=None,
) -> ExtractedPdf:
    """Async wrapper — runs the sync enricher in a worker thread."""
    import asyncio

    return await asyncio.to_thread(
        enrich_pdf_with_figures, path, extracted, progress_cb=progress_cb
    )


async def enrich_text_with_figures_async(
    path: Path | str,
    base_text: str,
    *,
    document_blocks=None,
    progress_cb=None,
) -> tuple[str, list[SheetGrid]]:
    import asyncio

    return await asyncio.to_thread(
        enrich_text_with_figures, path, base_text,
        document_blocks=document_blocks, progress_cb=progress_cb,
    )


async def enrich_office_document_with_figures_async(
    path: Path | str,
    base_text: str,
    *,
    document_blocks=None,
    progress_cb=None,
) -> OfficeFigureResult:
    import asyncio

    return await asyncio.to_thread(
        enrich_office_document_with_figures, path, base_text,
        document_blocks=document_blocks, progress_cb=progress_cb,
    )
