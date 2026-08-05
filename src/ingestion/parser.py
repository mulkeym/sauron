from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook


@dataclass
class Utterance:
    speaker: str
    text: str
    utterance_type: str  # "question", "statement"


@dataclass
class FigurePlacement:
    """An Office image occurrence anchored to its source document position.

    The relationship identifies the image bytes; ``figure_id`` identifies this
    particular occurrence.  Keeping those separate lets one repeated image be
    analysed once while retaining every placement and its local context.
    """

    figure_id: str
    relationship_id: str
    body_index: int
    section_path: list[str] = field(default_factory=list)
    caption: str = ""
    alt_text: str = ""
    previous_text: str = ""
    following_text: str = ""
    page: int | None = None
    bbox: tuple[float, float, float, float] | None = None


@dataclass
class DocumentBlock:
    """Ordered, source-aware content emitted by a document parser."""

    block_type: str  # "heading" | "paragraph" | "table" | "figure"
    body_index: int
    text: str = ""
    style: str = ""
    section_path: list[str] = field(default_factory=list)
    figure: FigurePlacement | None = None
    page: int | None = None
    bbox: tuple[float, float, float, float] | None = None


@dataclass
class ParsedDocument:
    filename: str
    doc_type: str  # "pdf", "docx", "pptx", "xlsx", "transcript"
    text: str
    metadata: dict = field(default_factory=dict)
    utterances: list[Utterance] = field(default_factory=list)
    blocks: list[DocumentBlock] = field(default_factory=list)


def parse_document(path: Path) -> ParsedDocument:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(path)
    elif suffix == ".docx":
        return _parse_docx(path)
    elif suffix == ".pptx":
        return _parse_pptx(path)
    elif suffix in (".xlsx", ".xlsm", ".xls", ".csv", ".tsv"):
        return _parse_spreadsheet(path)
    elif suffix == ".md":
        return _parse_markdown(path)
    elif suffix == ".txt":
        return _parse_transcript(path)
    elif suffix in (".conf", ".cfg", ".ini", ".yaml", ".yml", ".json",
                     ".xml", ".log", ".sh", ".bat", ".ps1", ".py",
                     ".html", ".htm", ".rtf"):
        return _parse_plaintext(path)
    else:
        # Try as plain text for any unrecognized extension
        return _parse_plaintext(path)


def _parse_plaintext(path: Path) -> ParsedDocument:
    """Parse any plain text file (config, log, script, etc.)."""
    raw = path.read_bytes()
    # Guard: refuse to ingest binary files as text. Without this, an
    # unhandled binary format (e.g. a legacy .xls routed here) gets decoded
    # into mojibake and silently embedded into the vector store.
    if b"\x00" in raw[:8192]:
        raise ValueError(
            f"{path.name}: appears to be a binary file, not text; refusing to ingest as plaintext"
        )
    text = raw.decode("utf-8", errors="replace")
    doc_type = path.suffix.lstrip(".") or "text"
    return ParsedDocument(filename=path.name, doc_type=doc_type, text=text)


def _parse_markdown(path: Path) -> ParsedDocument:
    text = path.read_text(encoding="utf-8")
    text = _strip_web_boilerplate(text)
    return ParsedDocument(filename=path.name, doc_type="markdown", text=text)


def _strip_web_boilerplate(text: str) -> str:
    """Strip web boilerplate from markdown documents.

    Works on any scraped website by removing lines that are structurally
    boilerplate (navigation links, social media, footers) while preserving
    actual content paragraphs. Generic enough for any document type.
    """
    import re

    lines = text.split("\n")
    kept = []

    for line in lines:
        stripped = line.strip()

        # Keep blank lines (preserve paragraph structure)
        if not stripped:
            kept.append(line)
            continue

        # Remove: lines that are purely markdown links with no meaningful text
        # e.g., "*   [News](https://...)" or "[](https://...)"
        plain_text = re.sub(r'\[([^\]]*)\]\([^)]*\)', r'\1', stripped)  # extract link text
        plain_text = re.sub(r'[*_\[\]()#|>!]', '', plain_text).strip()  # strip markdown formatting

        link_count = stripped.count('](')

        # Line is a navigation list item: "* [Link](url)" with minimal text
        if link_count >= 1 and len(plain_text) < 30 and stripped.startswith('*'):
            continue

        # Line is mostly links: 2+ links and plain text is minimal
        if link_count >= 2 and len(plain_text) < 20:
            continue

        # Line is an empty link: [](url) or [![Image](url)](url)
        if re.match(r'^\s*\[?\[?\]?\(', stripped) and len(plain_text) < 5:
            continue

        # Remove: social media links
        if any(s in stripped.lower() for s in ['facebook.com', 'instagram.com', 'linkedin.com',
               'youtube.com', 'twitter.com', '/#facebook', '/#x)', '/#email']):
            continue

        # Remove: common website boilerplate phrases
        if any(p in stripped.lower() for p in [
            'skip to main content', 'official websites use .gov',
            'secure .gov websites', 'how you know', 'share sensitive information',
            'official government organization', 'safely connected',
            'addtoany', 'thanks for sharing', 'previous next slideshow',
        ]):
            continue

        # Remove: sharing widgets
        if stripped in ('×', 'Share', '**Copy Link**', '---', 'Search', 'Search Search'):
            continue

        # Remove: image-only lines
        if re.match(r'^!\[.*\]\(.*\)$', stripped) or re.match(r'^\[!\[.*\]\(.*\)\]\(.*\)$', stripped):
            continue

        kept.append(line)

    # Collapse multiple consecutive blank lines
    result = re.sub(r'\n{4,}', '\n\n\n', "\n".join(kept)).strip()

    # Safety: if we removed >90%, something went wrong — keep original
    if len(result) < len(text) * 0.1 and len(text) > 500:
        return text

    return result


def _parse_pdf(path: Path) -> ParsedDocument:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages = []
        for page in reader.pages:
            page_text = page.extract_text() or ""
            if page_text.strip():
                pages.append(page_text)
        text = "\n".join(pages)
        return ParsedDocument(filename=path.name, doc_type="pdf", text=text)
    except Exception:
        from unstructured.partition.pdf import partition_pdf
        elements = partition_pdf(str(path))
        text = "\n".join(str(el) for el in elements)
        return ParsedDocument(filename=path.name, doc_type="pdf", text=text)


def _docx_para_text(para) -> str | None:
    """Format one paragraph; None if empty."""
    text = (para.text or "").strip()
    if not text:
        return None
    style = getattr(getattr(para, "style", None), "name", "") or ""
    if style.startswith("Heading"):
        return f"\n## {text}\n"
    return text


def _docx_table_text(table) -> str | None:
    """Flatten a Word table to readable prose (story-book layouts use tables heavily)."""
    rows_out: list[str] = []
    for row in table.rows:
        cells: list[str] = []
        seen: set[str] = set()
        for cell in row.cells:
            # python-docx repeats merged cells — de-dupe consecutive identical text
            ct = "\n".join(
                p.text.strip() for p in cell.paragraphs if p.text and p.text.strip()
            ).strip()
            if not ct or ct in seen:
                continue
            seen.add(ct)
            cells.append(ct)
        if cells:
            rows_out.append(" | ".join(cells) if len(cells) > 1 else cells[0])
    if not rows_out:
        return None
    return "\n".join(rows_out)


def _iter_docx_blocks(doc):
    """Yield paragraphs and tables in document-body order.

    ``Document.paragraphs`` skips table cell text entirely — children's books
    and many designed layouts put the story in tables, so ordered body walk
    is required for correct capture.
    """
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = doc.element.body
    for child in body.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, doc)
        elif isinstance(child, CT_Tbl):
            yield Table(child, doc)


def _docx_heading_level(style_name: str) -> int | None:
    if not style_name.startswith("Heading"):
        return None
    try:
        return max(1, int(style_name.split()[-1]))
    except (TypeError, ValueError):
        return 1


def _docx_image_refs(element) -> list[tuple[str, str]]:
    """Return ``(relationship id, alt text)`` pairs in XML order."""
    from docx.oxml.ns import qn

    refs: list[tuple[str, str]] = []
    for blip in element.xpath(".//a:blip"):
        rel_id = blip.get(qn("r:embed"))
        if not rel_id:
            continue
        alt_text = ""
        node = blip
        while node is not None:
            if node.tag in (qn("w:drawing"), qn("w:pict")):
                props = node.xpath(".//wp:docPr")
                if props:
                    prop = props[0]
                    alt_text = (
                        prop.get("descr") or prop.get("title") or ""
                    ).strip()
                break
            node = node.getparent()
        refs.append((rel_id, alt_text))
    return refs


def _docx_paragraph_units(para) -> list[tuple[str, str, str]]:
    """Split a paragraph into text/figure units without losing image anchors.

    Pictures normally occupy their own run.  Walking runs keeps their placement
    relative to surrounding text while retaining the paragraph style.
    """
    style = getattr(getattr(para, "style", None), "name", "") or ""
    if not _docx_image_refs(para._p):
        text = (para.text or "").strip()
        return [("text", text, style)] if text else []

    units: list[tuple[str, str, str]] = []
    buffered: list[str] = []
    emitted_refs: list[tuple[str, str]] = []

    def flush_text() -> None:
        text = "".join(buffered).strip()
        buffered.clear()
        if text:
            units.append(("text", text, style))

    for run in para.runs:
        if run.text:
            buffered.append(run.text)
        refs = _docx_image_refs(run._r)
        if refs:
            flush_text()
            for rel_id, alt_text in refs:
                units.append(("figure", rel_id, alt_text))
                emitted_refs.append((rel_id, alt_text))
    flush_text()

    # Older python-docx versions omit runs nested in hyperlinks and some other
    # smart tags. Preserve any drawing the paragraph XML saw but ``para.runs``
    # did not expose, anchored at the end of that paragraph as a safe fallback.
    remaining = list(emitted_refs)
    for ref in _docx_image_refs(para._p):
        if ref in remaining:
            remaining.remove(ref)
        else:
            units.append(("figure", ref[0], ref[1]))

    # Hyperlink text is not exposed consistently by older python-docx builds.
    # Preserve it as paragraph context when the run walk found only figures.
    if not any(kind == "text" for kind, _, _ in units):
        text = (para.text or "").strip()
        if text:
            units.insert(0, ("text", text, style))
    return units


def _nearest_text(blocks: list[DocumentBlock], start: int, direction: int) -> str:
    i = start + direction
    while 0 <= i < len(blocks):
        block = blocks[i]
        if (
            block.block_type != "figure"
            and not block.style.lower().startswith("caption")
            and block.text.strip()
        ):
            return block.text.strip()
        i += direction
    return ""


def render_document_blocks(
    blocks: list[DocumentBlock],
    figure_text: dict[str, str] | None = None,
) -> str:
    """Render ordered blocks, optionally inserting analysed figure prose."""
    figure_text = figure_text or {}
    rendered: list[str] = []
    for block in blocks:
        if block.block_type == "figure":
            if block.figure:
                text = (figure_text.get(block.figure.figure_id) or "").strip()
                if text:
                    rendered.append(text)
            continue
        text = block.text.strip()
        if not text:
            continue
        if block.block_type == "heading":
            rendered.append(f"## {text}")
        else:
            rendered.append(text)
    return "\n\n".join(rendered)


def _parse_docx(path: Path) -> ParsedDocument:
    doc = DocxDocument(str(path))
    blocks: list[DocumentBlock] = []
    seen_norm: set[str] = set()
    section_path: list[str] = []
    next_figure = 1

    def _add_text(text: str | None, block_type: str, style: str = "") -> None:
        nonlocal section_path
        if not text:
            return
        # Collapse whitespace for de-dupe (cover page often repeats title)
        key = " ".join(text.split())
        if key in seen_norm:
            return
        seen_norm.add(key)
        clean = text.strip()
        level = _docx_heading_level(style)
        if level is not None:
            section_path = section_path[:level - 1]
            section_path.append(clean)
            block_type = "heading"
        blocks.append(DocumentBlock(
            block_type=block_type,
            body_index=len(blocks),
            text=clean,
            style=style,
            section_path=list(section_path),
        ))

    def _add_figure(rel_id: str, alt_text: str = "") -> None:
        nonlocal next_figure
        placement = FigurePlacement(
            figure_id=f"fig-{next_figure:03d}",
            relationship_id=rel_id,
            body_index=len(blocks),
            section_path=list(section_path),
            alt_text=alt_text,
        )
        blocks.append(DocumentBlock(
            block_type="figure",
            body_index=len(blocks),
            section_path=list(section_path),
            figure=placement,
        ))
        next_figure += 1

    for block in _iter_docx_blocks(doc):
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        if isinstance(block, Paragraph):
            for kind, value, extra in _docx_paragraph_units(block):
                if kind == "figure":
                    _add_figure(value, extra)
                else:
                    _add_text(value, "paragraph", extra)
        elif isinstance(block, Table):
            table_text = _docx_table_text(block)
            _add_text(table_text, "table")
            # Images in cells are anchored to the table block.  They retain the
            # table text as local context even though Word has no useful page map.
            for rel_id, alt_text in _docx_image_refs(block._tbl):
                _add_figure(rel_id, alt_text)

    # Fill context after the full body is known so following captions are visible.
    for i, block in enumerate(blocks):
        if block.block_type != "figure" or not block.figure:
            continue
        previous = _nearest_text(blocks, i, -1)
        following = _nearest_text(blocks, i, 1)
        block.figure.previous_text = previous
        block.figure.following_text = following
        adjacent = []
        if i > 0:
            adjacent.append(blocks[i - 1])
        if i + 1 < len(blocks):
            adjacent.append(blocks[i + 1])
        for candidate in adjacent:
            if candidate.style.lower().startswith("caption") and candidate.text.strip():
                block.figure.caption = candidate.text.strip()
                break

    # Headers / footers (page numbers etc. — usually small, but keep searchable)
    for section in doc.sections:
        for hf in (section.header, section.footer):
            if hf is None:
                continue
            for para in hf.paragraphs:
                _add_text(_docx_para_text(para), "paragraph")
            for table in hf.tables:
                _add_text(_docx_table_text(table), "table")

    text = render_document_blocks(blocks)
    return ParsedDocument(filename=path.name, doc_type="docx", text=text, blocks=blocks)


def _pptx_shape_bbox(shape) -> tuple[float, float, float, float]:
    left = float(getattr(shape, "left", 0) or 0)
    top = float(getattr(shape, "top", 0) or 0)
    width = float(getattr(shape, "width", 0) or 0)
    height = float(getattr(shape, "height", 0) or 0)
    return left, top, left + width, top + height


def _pptx_image_ref(shape) -> tuple[str, str]:
    """Slide-local image relationship and authored alternative text."""
    try:
        from pptx.oxml.ns import qn

        blips = shape._element.xpath(".//a:blip")
        rel_id = blips[0].get(qn("r:embed")) if blips else ""
        props = shape._element.xpath(".//p:cNvPr")
        alt_text = ""
        if props:
            alt_text = (props[0].get("descr") or props[0].get("title") or "").strip()
        return rel_id or "", alt_text
    except Exception:
        return "", ""


def _pptx_table_text(shape) -> str:
    rows: list[str] = []
    try:
        for row in shape.table.rows:
            cells = [" ".join(cell.text.split()) for cell in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
    except Exception:
        return ""
    return "\n".join(rows)


def _pptx_chart_text(shape) -> str:
    """Extract chart categories and series values without needing a screenshot."""
    try:
        chart = shape.chart
    except Exception:
        return ""
    parts = ["[Chart]"]
    try:
        if chart.has_title and chart.chart_title.has_text_frame:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                parts.append(f"Title: {title}")
    except Exception:
        pass
    try:
        categories = []
        if len(chart.plots) > 0:
            categories = [str(c.label) for c in chart.plots[0].categories]
        if categories:
            parts.append("Categories: " + " | ".join(categories))
    except Exception:
        pass
    try:
        for series in chart.series:
            name = str(getattr(series, "name", "") or "Series")
            values = ["" if v is None else str(v) for v in series.values]
            parts.append(f"Series {name}: " + " | ".join(values))
    except Exception:
        pass
    parts.append("[/Chart]")
    return "\n".join(parts) if len(parts) > 2 else ""


def _iter_pptx_shapes(shapes):
    """Yield top-level and grouped shapes, keeping a stable visual order."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    ordered = sorted(
        list(shapes),
        key=lambda s: (
            float(getattr(s, "top", 0) or 0),
            float(getattr(s, "left", 0) or 0),
        ),
    )
    for shape in ordered:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_pptx_shapes(shape.shapes)
        else:
            yield shape


def _parse_pptx(path: Path) -> ParsedDocument:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(path))
    blocks: list[DocumentBlock] = []

    def add_text(
        text: str, block_type: str, slide_index: int,
        section_path: list[str], bbox=None,
    ) -> None:
        clean = (text or "").strip()
        if not clean:
            return
        blocks.append(DocumentBlock(
            block_type=block_type, body_index=len(blocks), text=clean,
            section_path=list(section_path), page=slide_index, bbox=bbox,
        ))

    for slide_index, slide in enumerate(presentation.slides):
        title_shape = slide.shapes.title
        title = ""
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title = (title_shape.text or "").strip()
        slide_label = f"Slide {slide_index + 1}"
        heading = f"{slide_label}: {title}" if title else slide_label
        section_path = [title or slide_label]
        add_text(heading, "heading", slide_index, section_path, (0, 0, 0, 0))

        slide_figure_no = 0
        for shape in _iter_pptx_shapes(slide.shapes):
            if title_shape is not None and shape._element is title_shape._element:
                continue
            bbox = _pptx_shape_bbox(shape)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                rel_id, alt_text = _pptx_image_ref(shape)
                if not rel_id:
                    continue
                slide_figure_no += 1
                placement = FigurePlacement(
                    figure_id=f"s{slide_index + 1}-fig-{slide_figure_no:03d}",
                    relationship_id=rel_id, body_index=len(blocks),
                    section_path=list(section_path), alt_text=alt_text,
                    page=slide_index, bbox=bbox,
                )
                blocks.append(DocumentBlock(
                    block_type="figure", body_index=len(blocks),
                    section_path=list(section_path), figure=placement,
                    page=slide_index, bbox=bbox,
                ))
                continue
            if getattr(shape, "has_table", False):
                add_text(_pptx_table_text(shape), "table", slide_index, section_path, bbox)
                continue
            if getattr(shape, "has_chart", False):
                add_text(_pptx_chart_text(shape), "chart", slide_index, section_path, bbox)
                continue
            if getattr(shape, "has_text_frame", False):
                add_text(shape.text, "paragraph", slide_index, section_path, bbox)

        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes and notes not in (str(slide_index + 1), title):
            add_text(
                f"Speaker notes:\n{notes}", "notes", slide_index,
                section_path, (0, float("inf"), 0, float("inf")),
            )

        # Attach captions and neighboring visual-order text on this slide.
        slide_indexes = [i for i, b in enumerate(blocks) if b.page == slide_index]
        for i in slide_indexes:
            block = blocks[i]
            if block.block_type != "figure" or not block.figure:
                continue
            prior = [blocks[j] for j in slide_indexes if j < i and blocks[j].text.strip()]
            after = [blocks[j] for j in slide_indexes if j > i and blocks[j].text.strip()]
            candidates = (prior[-2:] + after[:2])
            for candidate in candidates:
                if re.match(r"(?i)^\s*(figure|fig\.)\s*\d+", candidate.text):
                    block.figure.caption = candidate.text.strip()
                    break
            # A caption describes the picture, but it is not the surrounding
            # narrative context. Keep it in its dedicated field and walk past
            # it when selecting the before/after text sent to image analysis.
            def _is_caption(candidate: DocumentBlock) -> bool:
                return bool(
                    candidate.text.strip() == block.figure.caption
                    or re.match(r"(?i)^\s*(figure|fig\.)\s*\d+", candidate.text)
                )

            prior_context = [candidate for candidate in prior if not _is_caption(candidate)]
            after_context = [candidate for candidate in after if not _is_caption(candidate)]
            block.figure.previous_text = prior_context[-1].text if prior_context else ""
            block.figure.following_text = after_context[0].text if after_context else ""

    text = render_document_blocks(blocks)
    return ParsedDocument(filename=path.name, doc_type="pptx", text=text, blocks=blocks)


def _sniff_workbook_format(path: Path) -> str | None:
    """Detect a workbook's real format from its leading magic bytes.

    Extensions lie: some sources (e.g. OPM) serve modern OOXML workbooks under
    a legacy ``.xls`` name. Dispatching on content instead of extension is what
    lets those parse correctly rather than failing or producing mojibake.

    Returns "ooxml" (ZIP-based .xlsx/.xlsm), "ole" (legacy binary .xls), or
    None (not a recognized binary workbook — treat as delimited text).
    """
    with open(path, "rb") as f:
        head = f.read(8)
    if head[:4] == b"PK\x03\x04":
        return "ooxml"
    if head == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        return "ole"
    return None


def _parse_spreadsheet(path: Path) -> ParsedDocument:
    """Parse a spreadsheet to text, dispatching on actual content.

    openpyxl reads OOXML (.xlsx/.xlsm), xlrd reads legacy OLE .xls, and
    .csv/.tsv are plain delimited text. We sniff magic bytes first so a
    mislabeled file (e.g. OOXML content named .xls) is parsed by the right
    reader instead of being read as binary mojibake.
    """
    fmt = _sniff_workbook_format(path)
    if fmt == "ooxml":
        return _parse_xlsx(path)
    if fmt == "ole":
        return _parse_legacy_xls(path)
    # Not a binary workbook — delimited text (covers .csv/.tsv and any
    # text content mislabeled with a spreadsheet extension).
    return _parse_delimited(path)


def _parse_xlsx(path: Path) -> ParsedDocument:
    import io

    # Read via BytesIO so openpyxl dispatches on content, not the filename —
    # it otherwise refuses any path not ending in .xlsx/.xlsm, which breaks
    # the common case of OOXML content served under a .xls name.
    with open(path, "rb") as fh:
        wb = load_workbook(io.BytesIO(fh.read()), read_only=True)
    sheet_names = list(wb.sheetnames)
    parts = []
    for sheet_name in sheet_names:
        ws = wb[sheet_name]
        parts.append(f"Sheet: {sheet_name}")
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(h) if h else "" for h in rows[0]]
        parts.append(" | ".join(headers))
        for row in rows[1:]:
            parts.append(" | ".join(str(c) if c is not None else "" for c in row))
    wb.close()
    text = "\n".join(parts)
    return ParsedDocument(filename=path.name, doc_type="xlsx", text=text, metadata={"sheet_names": sheet_names})


def _parse_legacy_xls(path: Path) -> ParsedDocument:
    """Parse a legacy binary .xls workbook via xlrd."""
    import xlrd

    book = xlrd.open_workbook(str(path))
    sheet_names = book.sheet_names()
    parts = []
    for sheet in book.sheets():
        parts.append(f"Sheet: {sheet.name}")
        for r in range(sheet.nrows):
            cells = sheet.row_values(r)
            parts.append(" | ".join("" if c is None else str(c) for c in cells))
    text = "\n".join(parts)
    return ParsedDocument(filename=path.name, doc_type="xls", text=text, metadata={"sheet_names": sheet_names})


def _parse_delimited(path: Path) -> ParsedDocument:
    """Parse a .csv/.tsv file into ' | '-joined rows."""
    import csv

    # Guard: a binary file that slipped past the workbook sniff must not be
    # decoded into mojibake and ingested as "delimited text".
    if b"\x00" in path.read_bytes()[:8192]:
        raise ValueError(
            f"{path.name}: appears to be a binary file, not delimited text; refusing to ingest"
        )
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    parts = []
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        for row in csv.reader(f, delimiter=delimiter):
            parts.append(" | ".join(row))
    text = "\n".join(parts)
    return ParsedDocument(filename=path.name, doc_type=path.suffix.lstrip(".") or "csv", text=text)


def _parse_transcript(path: Path) -> ParsedDocument:
    raw = path.read_text(encoding="utf-8")
    lines = raw.strip().split("\n")
    metadata = {}
    utterances = []
    text_parts = []

    for line in lines:
        line = line.strip()
        if not line or line == "---":
            continue
        if ":" in line and not any(line.startswith(f"{name}:") for name in _extract_speaker_names(lines)):
            key, _, value = line.partition(":")
            if key.strip() in ("Meeting", "Date", "Location", "Attendees"):
                metadata[key.strip().lower()] = value.strip()
                text_parts.append(line)
                continue
        if ":" in line:
            speaker, _, text = line.partition(":")
            speaker = speaker.strip()
            text = text.strip()
            is_question = text.rstrip().endswith("?")
            utterances.append(Utterance(speaker=speaker, text=text, utterance_type="question" if is_question else "statement"))
            text_parts.append(line)
        else:
            text_parts.append(line)
    return ParsedDocument(filename=path.name, doc_type="transcript", text="\n".join(text_parts), metadata=metadata, utterances=utterances)


def _extract_speaker_names(lines: list[str]) -> set[str]:
    from collections import Counter
    names = Counter()
    for line in lines:
        if ":" in line:
            name = line.split(":")[0].strip()
            if name and len(name) < 40 and name not in ("Meeting", "Date", "Location", "Attendees"):
                names[name] += 1
    return {name for name, count in names.items() if count >= 1}
